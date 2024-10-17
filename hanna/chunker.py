import pandas as pd

from .credentials import ClientCredentials
# from langchain.text_splitter import RecursiveCharacterTextSplitter
from llama_index.core.node_parser import SentenceSplitter
from llama_index.core import Document
from pypdf import PdfReader
from openpyxl import load_workbook
import docx
import io
import logging
import htmltabletomd
from pptx import Presentation

logger = logging.getLogger(__name__)


class ChunkText(ClientCredentials):
    def __init__(self, size: int = 300):
        super(ChunkText).__init__()
        self.dir = "./_tmp/"

        self.__ppt_splitter = SentenceSplitter(chunk_size=6300)
        self.__rcts = SentenceSplitter(chunk_size=size)

        # RecursiveCharacterTextSplitter(chunk_overlap=0, chunk_size=size, length_function=len)

    def return_ext(self, filename: str):
        return filename.split(".")[-1]

    def split_table(self, table, chunk_size: int):
        for i in range(0, len(table), chunk_size):
            # Optionally drop columns with 'Unnamed' if they are irrelevant
            yield table[i: i + chunk_size].to_html()

    def df_to_markdown(self, df_chunk):
        return htmltabletomd.convert_table(df_chunk)

    def check_csv(self, filename) -> bool:
        df = pd.read_csv(self.dir + filename)

        logger.info(f"CSV ROWS {len(df)}")
        if 5 < len(df) < 501:
            return True
        else:
            logger.error("Too many rows!")
            return False

    def check_excel(self, filename: str) -> bool:
        check = False
        workbook = load_workbook(self.dir + filename)
        sheet_count = len(workbook.sheetnames)

        if sheet_count <= 3:
            for sheet in workbook.sheetnames:
                logger.info(f"SHEET NAME: {sheet}")
                row_count = workbook[sheet].max_row
                if 5 < row_count < 502:
                    check = True
                else:
                    check = False
                    logger.error("Too many rows!")
                    break

            logger.info(f"The Excel file '{filename}' has {sheet_count} sheets.")

        return check

    def check_ppt(self, filename: str) -> bool:
        prs = Presentation(self.dir + filename)
        logger.info(f"TOTAL SLIDES: {len(prs.slides)}")

        if len(prs.slides) > 100:
            logger.error("Too many slides, MAX 40!!!")
            return False

        return True

    def process_ppt(self, filename: str) -> bool | list[str]:
        tmp = []
        name = filename.split('.')[0]
        logger.info("READING PPTX...")
        prs = Presentation(self.dir + filename)

        for slide_number, slide in enumerate(prs.slides):

            for shape in slide.shapes:
                title = f"FILE: {name} \nTitle: {shape.title.text}" if hasattr(shape, "title") else "FILE: " + name

                if hasattr(shape, "text"):
                    tmp.append(f"{title} \nSlide number: {slide_number + 1} \nDescription: {shape.text}")

        return tmp

    def process_excel(self, filename: str):
        tables = []
        df = pd.read_excel(self.dir + filename, sheet_name=None)
        for sheet_name, dtf in df.items():
            dtf.dropna(how='all', inplace=True)

            for data in self.split_table(dtf, 50):
                mk = self.df_to_markdown(data)
                tables.append(mk)

        return tables

    def process_table(self, file_name: str):
        df = pd.read_excel(file_name, sheet_name=None)
        # res = htmltabletomd.convert_table()

        for sheet_name, dtf in df.items():
            dtf.dropna(how='all', inplace=True)

            for data in self.split_table(dtf, 50):
                mk = self.df_to_markdown(data)
                yield mk

    def chunk_document(self, path: str) -> list | None:
        corpus = ""
        file = path.split('.')

        if file[-1] == "docx":
            document = docx.Document(self.dir + path)
            print(f"\nReading DOCX [{file}...")

            for x, page in enumerate(document.paragraphs):
                # print(f"PAGE {x + 1}")
                corpus += str(page.text).lower() + " "

        elif file[-1] == "txt":
            file_txt = io.open(self.dir + path, encoding="utf-8")
            corpus = file_txt.read().lower()
            print(f"\nReading TXT [{file_txt.name}]...")
            file_txt.close()

        elif file[-1] == "pdf":
            reader = PdfReader(self.dir + path)
            for page in reader.pages:
                corpus += page.extract_text()

        elif file[-1] == "pptx":
            return self.process_ppt(path)

        elif file[-1] == "xlsx" or file[-1] == "csv":

            return None

        else:
            return None

        pre_filter = corpus.strip().replace("\n", "").lower()

        split_text = self.__rcts.split_text(pre_filter)

        return split_text

    def chunk_corpus(self, text: str):

        # split_text = self.__rcts.split_text(text)
        chunks = self.__rcts.get_nodes_from_documents([Document(text=text)])
        nodes = [chunk.text for chunk in chunks]

        return nodes

