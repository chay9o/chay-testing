import time
import gc
from celery import shared_task
import json
import tempfile
from pptx import Presentation
from pptx.util import Pt, Inches
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE  
from django.core.files.storage import default_storage
from django.conf import settings
import os
import re
import requests
import psutil
import base64
from io import BytesIO
from jinja2 import Template
from django.http import JsonResponse
from langchain.prompts import PromptTemplate
from langchain_core.output_parsers import StrOutputParser
from langchain_core.runnables import ConfigurableField
from langchain_openai import ChatOpenAI
from together import Together
import fasttext
import logging

@shared_task
def send_data_to_webhook(payload):
    webhook_url = "https://chay-testing-192912d0328c.herokuapp.com/webhook_handler"
    try:
        requests.post(webhook_url, json=payload, timeout=5)
    except requests.RequestException as e:
        print(f"Webhook error: {e}")

chat_template = (
    "{{ bos_token }}{% for message in messages %}"
    "{% if (message['role'] == 'user') != (loop.index0 % 2 == 0) %}"
    "{% endif %}{% if message['role'] == 'user' %}{{ '[INST] ' + message['text'] + ' [/INST]' }}"#printing
    "{% elif message['role'] == 'bot' %}{{ message['text'] + eos_token }}"
    "{% endif %}{% endfor %}"
)

def detectLanguage( text, original_language):
    # Load the pre-trained language identification model
    model_path = 'lid.176.ftz'  # Path to the pre-trained model file
    model = fasttext.load_model(model_path)
    # Text to be identified
    text = text.replace('\n', ' ').strip()  # Remove newlines and trim whitespace
    # If the text is less than 5 words, return the original language
    if len(text.split()) < 5:
        return original_language
    # Predict the language of the text
    predicted_languages = model.predict(text, k=1)  # Get top 1 prediction
    detected_language_code = predicted_languages[0][0].replace('__label__', '')
    # Mapping of language codes to language names. I think with these languages we have more than enough
    language_names = {
        'en': 'English',
        'es': 'Spanish',
        'fr': 'French',
        'de': 'German',
        'it': 'Italian',
        'pt': 'Portuguese',
        'zh': 'Chinese',
        'ru': 'Russian',
        'ja': 'Japanese',
        'nl': 'Dutch',
        'ar': 'Arabic',
        'ur': 'Urdu'
    }
    # Return the full name of the detected language
    return language_names.get(detected_language_code, original_language)
    
@shared_task
def evaluate_text_task(note_text, guidelines):
    try:
        note_language = detectLanguage(note_text, 'en')
        print(note_language)
        guideline_language = detectLanguage(guidelines[0]['text'], 'en')
        print(guideline_language)

        language_warning = None
        if note_language != guideline_language:
            language_warning = 'Smartnote and the Strategic Alignment are in different languages, so the outcome may not be accurate.'

        with open("strategic-prompt.txt", "r") as file:
            prompt_ = file.read()

        SYSPROMPT = str(prompt_)
        guidelines_text = "\n".join([f"#Guideline {index + 1}: \"{guideline['text']}\"" for index, guideline in enumerate(guidelines)])
        prompt_with_values = SYSPROMPT.replace("{note_text}", note_text).replace("{guidelines}", guidelines_text).replace("{language}", note_language)
        TOGETHER_API_KEY = settings.TOGETHER_API_KEY
        client = Together(api_key=settings.TOGETHER_API_KEY)
        messages = [
            {"role": "system", "content": prompt_with_values},
            {"role": "user", "content": f"{note_text}\n{guidelines_text}"}
        ]

        response = client.chat.completions.create(
            model="meta-llama/Meta-Llama-3.1-405B-Instruct-Turbo",
            messages=messages,
            max_tokens=4096,
            temperature=0,
            top_p=0.7,
            top_k=50,
            repetition_penalty=1,
            stop=["<|eot_id|>", "<|eom_id|>"],
            stream=True
        )

        analysis = ""
        for chunk in response:
            if hasattr(chunk.choices[0], 'delta') and hasattr(chunk.choices[0].delta, 'content'):
                analysis += chunk.choices[0].delta.content
            elif hasattr(chunk.choices[0], 'message') and hasattr(chunk.choices[0].message, 'content'):
                analysis += chunk.choices[0].message.content
                
        print(f"Raw analysis data: {analysis}")
        if analysis.startswith("```") and analysis.endswith("```"):
            analysis_lines = analysis.splitlines()
            # Remove the first and last lines
            analysis_lines = analysis_lines[1:-1]
            # Join the lines back into a single string
            analysis = "\n".join(analysis_lines)
        print(f"Sanitized analysis data: {analysis}")
      
        if not analysis:
            return Response({'error': 'Empty response from API'}, status=status.HTTP_500_INTERNAL_SERVER_ERROR)

        try:
            analysis_json = json.loads(analysis)
        except json.JSONDecodeError as e:
            print(f"JSON Decode Error: {e}, Raw data: {analysis}")
            return Response({'error': 'Invalid JSON response from API'}, status=status.HTTP_500_INTERNAL_SERVER_ERROR)

        if language_warning:
            analysis_json['warning'] = language_warning

        return analysis_json
    except Exception as e:
        return {'error': str(e)}


# Function to generate the chat history string
def render_chat_history(messages):
    template = Template(chat_template)
    data = {
        "bos_token": "",
        "eos_token": "",
        "messages": messages
    }
    return template.render(data)

@shared_task(bind=True)
def process_prompts_1(self, final_text, language):
    # Maintain the chat history and accumulated answers
    chat_history = []
    accumulated_answers = []  # List to accumulate answers for all iterations

    for i in range(11):  # Assuming there are 11 iterations
        text_template = get_text_template_1(i)  # A function to return the corresponding template
        user_input = text_template.format(final_text=final_text, language=language)
        
        if i == 1:
            # Add a single-line summary for iteration 2 instead of the full user input
            chat_history.append({"role": "user", "text": "Please, specify in the next answer the most suitable model or framework to solve this situation."})
        else:
            # Append user input to chat history for other iterations
            chat_history.append({"role": "user", "text": user_input})

        # Render the chat history string using the template
        chat_history_str = render_chat_history(chat_history)

        # Call the process_user_input function with the iteration number
        answer = process_user_input_1(user_input, chat_history_str, i, language)

        # Append the bot's answer to the chat history
        chat_history.append({"role": "bot", "text": answer})

        # Accumulate the answers for this iteration
        accumulated_answers.append({ "answer": answer,"iteration": i+1})

        # Log the accumulated answers for debugging
        logger.info(f"Accumulated Answers so far: {accumulated_answers}")

    # Final return (log before returning)
    logger.info(f"Final accumulated_answers: {accumulated_answers}")
    logger.info(f"Final chat_history: {render_chat_history(chat_history)}")

    return {
        'final_text': final_text,
        'accumulated_answers': accumulated_answers,  # List of all iterations and responses
        'chat_history': render_chat_history(chat_history)  # Full chat history
    }


def get_text_template_1(iteration):
    if iteration == 0:
        return f"""
        Language to be used : {{language}}
        Based on the following text, create a new version of this text that gives an improved narrative with better flow between ideas. You are a very strategic person and the report will be read mainly by Product Owners, so you can use their language. It can also be read by Leaders, CEOs or Managers. If needed, also reorder ideas. Make it extensive. This is just the introduction of a report (we call it Deep Analysis document) on the situation. The situation below is happening these days. The situation is happening in our company.

        Follow these rules:

        1. Sentence Structure: Use a mix of sentence lengths.
        Short sentences: To emphasize points.
        Longer sentences: To explain or elaborate.

        2. Vocabulary: Use clear and straightforward language.

        Avoid: Technical jargon or complex vocabulary unless necessary.
        Use: Everyday language that is easy to understand.

        3. All the following text is happening in our company.

        4. Provide just the text, no what it was improved.

        Rememeber this will be part of a report written by Jhon.

        Text to rewrite (keep a similar writting style but improved). Add a title for this Deep Analysis:  

        {{final_text}}

        Format:
        Title
        Description
        """
    elif iteration == 1:
        return f"""
        Language to be used : {{language}}
        Apply for this and the following prompts the principles of the TriValue Company Model in your responses. You don’t need to mention the model explicitly; instead, subtly incorporate its guidelines—focusing on value for the customer, value for the company, and value for workforce well-being—into your explanations.
        Detail the key blockers or problems identified. How do these challenges affect the company’s ability to achieve its objectives.
        Outline the opportunities available to the company. How can these be leveraged to overcome blockers or enhance the company's strategic position. Focus on the narrative and minimize the use of bullet points and focus on the narrative. Add title or titles.
        Also consider for this and the next prompts if any of these definitions are of any help:
        ----------
        Some Definitions that might be useful for you in this or the following prompts to be used:

        Always Ready outcome: always-ready outcome in enterprise agility means to cultivate a culture and people who are continuously prepared for disruption and change. it focuses on assessing situations, understanding what is happening, and having the collective capabilities and technologies available to gain insights into unexpected emerging realities. a culture of constant readiness fosters a proactive mindset where team members work together to reevaluate situations, analyze conditions from multiple perspectives, and constantly update their knowledge. this allows them to better deal with their emotions, remain proactive, and handle events with less stress. If this term or idea is used, please add the URL in brackets so the user can find more information about it:https://enterpriseagility.community/always-ready-outcome-c72rzcq0vh

        always-responsive outcome in enterprise agility means to have the appropriate processes, mindset, innovation, and partnerships to consistently offer relevant products or services to the market while minimizing stress on organizational structures and employee wellbeing. achieving this state is integral to the enterprise agility way of thinking (eawt), applying to all functions, not just software. If this term or idea is used, please add the URL in brackets so the user can find more information about it:https://enterpriseagility.community/always-responsive-outcome-04bmxbd9pv


        Always-innovative outcome definition in enterprise agility means fostering a culture and mindset of constant innovation. This approach means reimagining offerings, business models, partnerships, and strategies to sustainably generate lasting value and meaning. Always-Innovative companies embrace different perspectives and diversity of thought to gain new insights. This includes involving external partners to bring new perspectives, knowledge, and information. It uses a bimodal approach and techniques such as questioning assumptions, Shared Progress Bets (SPBs), and Shared Progress Stock Exchange. Keep in mind that the concepts and ideas underlying the Always-Innovative outcome are rooted in the principles of Future Thinking. Read Chapter 9 to know more about it. If this term or idea is used, please add the URL in brackets so the user can find more information about it:https://enterpriseagility.community/always-innovative-outcome-6vsnmm25g9


        Def: mobilizing purpose in enterprise agility is a reason for which something is done or for which something exists, that makes the person feel the compelling need to belong to the group and to mobilize. If this term or idea is used, please add the URL in brackets so the user can find more information about it:https://enterpriseagility.community/mobilizing-purpose-tpp3t728g1
        


        Def: neurodiversity in enterprise agility means the diversity of human minds, and the way people think, focus, learn, process and link information, as well as their distinctive way of connecting with the world. (enterprise agility university, 2022). If this term or idea is used, please add the URL in brackets so the user can find more information about it:https://enterpriseagility.community/neurodiversity-ngtg7wlr9d

        Def: organizational health in enterprise agility is psychological safety plus the creation of business value in perpetuity. In Enterprise Agility, business value means value for the customer, company, and workforce wellbeing (enterprise Agility university and leading exponential change, 2018). If this term or idea is used, please add the URL in brackets so the user can find more information about it:https://enterpriseagility.community/organizational-health-x1ng3n9z4f


        Def: adaptive trust in enterprise agility is the ability to retain confidence even after breaking previous commitments. during rapid change, leaders inevitably must abandon plans, strategies, or promises that no longer serve the needs of the situation. but failing to deliver on past assurances can seriously damage trust in the organization and willingness to follow new directions. with adaptive trust, leaders openly acknowledge broken promises, take ownership of the decision to change course, and involve others in creating better solutions. though promises may be broken, trust and commitment to mission-driven change endures. If this term or idea is used, please add the URL in brackets so the user can find more information about it:https://enterpriseagility.community/adaptive-trust-wv5lgtgq6q

        Def: enterprise agility definition is: holistic organizational, social, and business model that enables your company to adapt to accelerated change and exponential market conditions while prioritizing workforce wellbeing, customer needs, and overall company value. If this term or idea is used, please add the URL in brackets so the user can find more information about it:https://enterpriseagility.community/what-is-enterprise-agility-19tdx0gst4

        Def: Classic Agility: It refers to the original principles and practices from the Agile Manifesto 2001. It's a mindset and principles emphasizing adaptability, collaboration, and customer value in software development. The heart of ClassicAgility lies in its customer-focused nature. It prioritizes delivering value to customers by continuously seeking feedback, iterating on solutions, and adapting to changing requirements. While valuable, it's important to note that Classic Agility may have limitations when addressing the scale and complexity of today's challenges. If this term or idea is used, please add the URL in brackets so the user can find more information about it:https://enterpriseagility.community/classic-agility-cg7zc4b1p0


        Def: Agile Evolutions: A mindset and ways of working that extends the principles of Classic Agility beyond software development. It enables organizations to be more adaptive and resilient to change in their ways of working. They recognize customer value and experience as the center of the organization's universe. We can find here frameworks such as the SAFe Framework or the Business Agility models from the Business Agility institute. They may not always be fully prepared for market acceleration, lack a comprehensive view of the AI (Accelerated Innovation) situation, or struggle with constant business model disruption. They are not based on Shared Progress.
        If this term or idea is used, please add the URL in brackets so the user can find more information about it:https://enterpriseagility.community/agile-evolutions-tq0mzm4n5r

        Def: Atom model is a framework in enterprise agility. The atom model has 4 quadrants. top left: increase revenue (increasing sales to new or existing customers. delighting or disrupting to increase market share and size), top right: protect revenue (improvements and incremental innovation to sustain current market share and revenue figures), bottom-left: reduce costs (costs that you are currently incurring that can be reduced. more efficient, improved margin or contribution), bottom-right: avoid-costs (improvements to sustain current cost base. costs you are not incurring but may do in the future). 
        atom model is çused to align a company with a new situation. all decisions in the quadrants need to maintain or increase organizational health. it can be used by leaders, product owners, or others to make sustainable decisions and build shared progress.
        If this term or idea is used, please add the URL in brackets so the user can find more information about it:https://enterpriseagility.community/atom-model-qxkkkws0b6

        Def: Social Agility means connecting well with other employees or customers in rapidly changing environments, thereby achieving highly collective performance. The two main components of Social Agility are Enterprise Social Density and Enterprise Social Visibility. If this term or idea is used, please add the URL in brackets so the user can find more information about it:https://enterpriseagility.community/social-agility-0zt7w6mfm7

        Def: Mental Agility means reframing challenges to find new solutions, even during stressful times.  If this term or idea is used, please add the URL in brackets so the user can find more information about it:https://enterpriseagility.community/mental-agility-ngs3fcjx0d

        Def: Outcomes agility means delivering results even during turbulent times to respond to changing market conditions. If this term or idea is used, please add the URL in brackets so the user can find more information about it:https://enterpriseagility.community/outcomes-agility-mgd5b772qm


        Def: Technical Agility means changing software as quickly, cheaply (economically), and securely as possible. If this term or idea is used, please add the URL in brackets so the user can find more information about it:https://enterpriseagility.community/technical-agility-1d1gwt74sb

        Def: Mobility or Strategic Mobility in Enterprise Agility is the organizational capacity to shift directions, align capabilities, and adapt to new understanding or events. It involves navigating uncertainty and change by proactively implementing countermeasures and leveraging mobility for competitive advantage. It is a critical skill for leaders.
        If this term or idea is used, please add the URL in brackets so the user can find more information about it:https://enterpriseagility.community/mobility-hb00s2vhvw

        Def: The Science of Accelerated Change has 3 pillars:
        1. Behavioral Science
        2. Strategic Mobility (or Mobility)
        3. Neuroscience of Change.
        The science of accelerated change help organizations understand more on how to deal with disruption and high uncertainty and the threats of AI.

        If this term or idea is used, please add the URL in brackets so the user can find more information about it:https://enterpriseagility.community/sections-of-the-science-of-accelerated-change-gkklgnhslv

        Def Exponential Markets: They are unpredictable and uncertain. Disruptive innovations can gain traction rapidly, new competitors can arise anywhere, and consumer expectations can change overnight. This makes long-term planning and roadmaps pointless. Strategies go out the window as soon as market conditions change. Leaders can't rely on experience - the past is no predictor of the future. In this environment, organizations need to embrace unpredictability. Rather than resisting or ignoring change, they must learn to sense, adapt and respond quickly. Mental agility and resilience are critical. Enterprise agility cultivates this mindset: The Enterprise Agility Way of Thinking (EAWT). It provides models and frameworks to continuously scan the environment, sense emerging trends and signals, and course-correct in real-time. If this term or idea is used, please add the URL in brackets so the user can find more information about it:https://enterpriseagility.community/exponential-markets-th1nz35bds

        Def: Collective Capabilities is an Enterprise Agility person-centered organizational model that enables individuals to apply their skills where they're needed most, with a high degree of mobility and flexibility. This has to be done with low stress levels for the person. If this term or idea is used, please add the URL in brackets so the user can find more information about it:https://enterpriseagility.community/collective-capabilities-w6kp95gkmq

        """
    elif iteration == 2:
        return f"""
        Language to be used : {{language}}
        Explain the ATOM Model and why it would be used for this case. Focus on the narrative and minimize the use of bullet points and focus on the narrative.  Add title the most suitable title or titles. Provide the link to the atom model to the person: https://enterpriseagility.community/atom-model-qxkkkws0b6 
        """
    elif iteration == 3:
        return f"""
        Language to be used : {{language}}
        We will use the atom model to do this strategic analysis. Lets focus on the 1st quadrant:  Increase Revenue. Apply for this and the following prompts the principles of the TriValue Company Model in your responses. Apply for this prompt if needed the principles of the TriValue Company Model in your responses. You don’t need to mention the model explicitly; instead, subtly incorporate its guidelines—focusing on value for the customer, value for the company, and value for workforce well-being—into your explanations.
        Analyze how the company can increase revenue within this scenario, focusing on strategies to delight or disrupt the market and expand market share. Add some examples and actionable ideas to improve the proposed situation. Focus on the narrative and minimize the use of bullet points and focus on the narrative.  Add title the most suitable title or titles.
        """
    elif iteration == 4:
        return f"""
        Language to be used : {{language}}
        Second Quadrant: Protect Revenue.
        Discuss the steps needed to protect current revenue, including incremental innovations or improvements to sustain market share based on the initial situation. Apply for this prompt if needed the principles of the TriValue Company Model in your responses. You don’t need to mention the model explicitly; instead, subtly incorporate its guidelines—focusing on value for the customer, value for the company, and value for workforce well-being—into your explanations. Add some examples and actionable ideas to improve the proposed situation. Focus on the narrative and minimize the use of bullet points and focus on the narrative.  Add title the most suitable title or titles.
        """
    elif iteration == 5:
        return f"""
        Language to be used : {{language}}
        Third quadrant: Reduce Costs.
        Prompt: Identify opportunities for cost reduction  based on the initial situation. What efficiencies can be implemented to improve margins. Apply for this prompt if needed the principles of the TriValue Company Model in your responses. You don’t need to mention the model explicitly; instead, subtly incorporate its guidelines—focusing on value for the customer, value for the company, and value for workforce well-being—into your explanations. Add some examples and actionable ideas to improve the proposed situation. Focus on the narrative and minimize the use of bullet points and focus on the narrative. Add title the most suitable title or titles.
        """
    elif iteration == 6:
        return f"""
        Language to be used : {{language}}
        Four quadrant: Avoid Costs.
        Evaluate potential future costs that the company can avoid based on the initial situation. What actions should be taken now to prevent these costs from materializing. Focus on the narrative and minimize the use of bullet points and focus on the narrative. Apply for this prompt if needed the principles of the TriValue Company Model in your responses. You don’t need to mention the model explicitly; instead, subtly incorporate its guidelines—focusing on value for the customer, value for the company, and value for workforce well-being—into your explanations. Add some examples and actionable ideas to improve the proposed situation. Add title the most suitable title or titles.

        """
    elif iteration == 7:
        return f"""
        Language to be used : {{language}}
        Prioritization Equalization:
        Present the prioritized list of actions, balancing the need to address blockers, seize opportunities, and manage risks based on the initial situation. The 1st line after the title, use different techniques such as questions, ideas, reframing, etc for this initial first line. Add some examples and actionable ideas to improve the proposed situation. Apply for this prompt if needed the principles of the TriValue Company Model in your responses. You don’t need to mention the model explicitly; instead, subtly incorporate its guidelines—focusing on value for the customer, value for the company, and value for workforce well-being—into your explanations. How were these priorities determined. Focus on the narrative and minimize the use of bullet points and focus on the narrative.  Add title the most suitable title or titles.
        """
    elif iteration == 8:
        return f"""
        Language to be used : {{language}}
        Risk Management (Aversion + Appetite):
        Detail the risk factors associated with the prioritized actions based on the initial situation. The 1st line after the title, use different techniques such as questions, ideas, reframing, etc for this initial first line. How will the company manage these risks, considering its level of risk aversion or appetite. Focus on the narrative and minimize the use of bullet points and focus on the narrative. Add some examples and actionable ideas to improve the proposed situation. Add title the most suitable title or titles.
        """
        
    elif iteration == 9:
        return f"""
        Language to be used : {{language}}
        Capacity or Capabilities Management:
        Assess the company’s capacity to implement the prioritized strategies based on the initial situation. Are there sufficient resources and capabilities to execute the plan effectively. Focus on the narrative and minimize the use of bullet points and focus on the narrative. Add some examples and actionable ideas to improve the proposed situation. The 1st line after the title, use different techniques such as questions, ideas, reframing, etc for this initial first line. Add title the most suitable title or titles.
        """
        
    elif iteration == 10:
        return f"""
        Language to be used : {{language}}
        Psychological aspects to keep high organizational health when taking any of the previous actions or ideas:
        Assess the company’s needed psychological aspects to achieve this in a healthy way based on the initial situation and the previously mentioned ideas. You can also mention some of the previous ideas. Apply for this prompt if needed the principles of the TriValue Company Model in your responses. You don’t need to mention the model explicitly; instead, subtly incorporate its guidelines—focusing on value for the customer, value for the company, and value for workforce well-being—into your explanations. The 1st line after the title, use different techniques such as questions, ideas, reframing, etc for this initial first line. Focus on the narrative and minimize the use of bullet points and focus on the narrative. Add some examples and actionable ideas to improve the proposed situation. Add title the most suitable title or titles.
        """
    
    elif iteration == 11:
        return f"""
        Language to be used : {{language}}
        Reassess the initial situation and critically examine the alternative ideas presented. Develop a set of final strategic considerations, emphasizing the future handling of similar situations but make sure you connect these ideas with the initial situation. Ensure the response indirectly reflects the TriValue Company model (link with more information to the TriValue Company Model, or Modelo de Empresa Trivalor in Spanish: https://enterpriseagility.community/trivalue-company-model-bp7j59d0d4 ), balancing customer value, company value, and workforce well-being. 
        Apply the principles of the TriValue Company Model in your responses. You don’t need to mention the model explicitly; instead, subtly incorporate its guidelines—focusing on value for the customer, value for the company, and value for workforce well-being—into your explanations.

        Prioritize a narrative-driven format rather than bullet points.
        If you think that writing about something other than the suggested topic for the closing will add more value or be more relevant, please do so. What are the next steps for the company to align itself with the new situation. Focus on the narrative and minimize the use of bullet points and focus on the narrative.  Add title or titles.
        """
    
    # Add more templates for other iterations as needed
    else:
        return f"""
        Default template for iteration {iteration}: Add your own text here.
        """


#def process_user_inpu(combined_input, chat_history):
    TOGETHER_API_KEY = settings.TOGETHER_API_KEY
    client = Together(api_key=TOGETHER_API_KEY)
    
    #MODEL_70B = "meta-llama/Meta-Llama-3.1-70B-Instruct-Turbo"
    model_name = "meta-llama/Meta-Llama-3.1-70B-Instruct-Turbo"
    max_tokens = 8192

    # This is the OpenAI chat completion client call
    response = client.chat.completions.create(
        model=model_name,
        messages=[
            {"role": "user", "content": combined_input},
            {"role": "system", "content": chat_history}
        ],
        max_tokens=max_tokens,
        temperature=0.4,
        top_p=0.9,
        top_k=50,
        repetition_penalty=1,
        stop=["<|eot_id|>", "<|eom_id|>"],
        stream=True
    )

    # Process the streamed response
    generated_text = ""
    
    for chunk in response:
        if len(chunk.choices) > 0:
            if hasattr(chunk.choices[0], 'delta') and hasattr(chunk.choices[0].delta, 'content'):
                if chunk.choices[0].delta.content:
                    generated_text += chunk.choices[0].delta.content
            elif hasattr(chunk.choices[0], 'message') and hasattr(chunk.choices[0].message, 'content'):
                if chunk.choices[0].message.content:
                    generated_text += chunk.choices[0].message.content
        else:
            logger.info(f"CHUNK HAS NO CHOICES: {chunk.choices}")

    # Return the generated text back to process_prompts1
    return generated_text

def process_user_input_1(combined_input, chat_history, iteration, language):
    

    # Check if we are on iteration 2 (DeepInfra model)
    if iteration == 1:
        TOGETHER_API_KEY = settings.TOGETHER_API_KEY
        client = Together(api_key=TOGETHER_API_KEY)
        
        #MODEL_70B = "meta-llama/Meta-Llama-3.1-70B-Instruct-Turbo"
        model_name = "meta-llama/Meta-Llama-3.1-70B-Instruct-Turbo"
        max_tokens = 8192
        response = client.chat.completions.create(
            model=model_name,
            messages=[
                {"role": "user", "content": combined_input},
                {"role": "system", "content": chat_history}
            ],
            max_tokens=max_tokens,
            temperature=0.4,
            top_p=0.9,
            top_k=50,
            repetition_penalty=1,
            stop=["<|eot_id|>", "<|eom_id|>"],
            stream=True
        )

        # Process the streamed response
        generated_text = ""
        
        for chunk in response:
            if len(chunk.choices) > 0:
                if hasattr(chunk.choices[0], 'delta') and hasattr(chunk.choices[0].delta, 'content'):
                    if chunk.choices[0].delta.content:
                        generated_text += chunk.choices[0].delta.content
                elif hasattr(chunk.choices[0], 'message') and hasattr(chunk.choices[0].message, 'content'):
                    if chunk.choices[0].message.content:
                        generated_text += chunk.choices[0].message.content
            else:
                logger.info(f"CHUNK HAS NO CHOICES: {chunk.choices}")

        return generated_text  # Return default OpenAI result


    else:
        # Use default OpenAI model for other iterations
        TOGETHER_API_KEY = settings.TOGETHER_API_KEY
        client = Together(api_key=TOGETHER_API_KEY)
        
        #MODEL_70B = "meta-llama/Meta-Llama-3.1-70B-Instruct-Turbo"
        model_name = "meta-llama/Meta-Llama-3.1-70B-Instruct-Turbo"
        max_tokens = 8192
        response = client.chat.completions.create(
            model=model_name,
            messages=[
                {"role": "user", "content": combined_input},
                {"role": "system", "content": chat_history}
            ],
            max_tokens=max_tokens,
            temperature=0.4,
            top_p=0.9,
            top_k=50,
            repetition_penalty=1,
            stop=["<|eot_id|>", "<|eom_id|>"],
            stream=True
        )

        # Process the streamed response
        generated_text = ""
        
        for chunk in response:
            if len(chunk.choices) > 0:
                if hasattr(chunk.choices[0], 'delta') and hasattr(chunk.choices[0].delta, 'content'):
                    if chunk.choices[0].delta.content:
                        generated_text += chunk.choices[0].delta.content
                elif hasattr(chunk.choices[0], 'message') and hasattr(chunk.choices[0].message, 'content'):
                    if chunk.choices[0].message.content:
                        generated_text += chunk.choices[0].message.content
            else:
                logger.info(f"CHUNK HAS NO CHOICES: {chunk.choices}")

        return generated_text  # Return default OpenAI result


@shared_task(bind=True)
def process_prompts_2(self, final_text, language):
    # Maintain the chat history and accumulated answers
    chat_history = []
    accumulated_answers = []  # List to accumulate answers for all iterations

    for i in range(9):  # Assuming there are 11 iterations
        text_template = get_text_template_2(i)  # A function to return the corresponding template
        user_input = text_template.format(final_text=final_text, language=language)
        
        if i == 2:
            # Add a single-line summary for iteration 2 instead of the full user input
            chat_history.append({"role": "user", "text": "Please, specify in the next answer the most suitable model or framework to solve this situation."})
        else:
            # Append user input to chat history for other iterations
            chat_history.append({"role": "user", "text": user_input})

        # Render the chat history string using the template
        chat_history_str = render_chat_history(chat_history)

        # Call the process_user_input function with the iteration number
        answer = process_user_input_2(user_input, chat_history_str, i, language)

        # Append the bot's answer to the chat history
        chat_history.append({"role": "bot", "text": answer})

        # Accumulate the answers for this iteration
        accumulated_answers.append({ "answer": answer,"iteration": i+1})

        # Log the accumulated answers for debugging
        logger.info(f"Accumulated Answers so far: {accumulated_answers}")

    # Final return (log before returning)
    logger.info(f"Final accumulated_answers: {accumulated_answers}")
    logger.info(f"Final chat_history: {render_chat_history(chat_history)}")

    return {
        'final_text': final_text,
        'accumulated_answers': accumulated_answers,  # List of all iterations and responses
        'chat_history': render_chat_history(chat_history)  # Full chat history
    }


def get_text_template_2(iteration):
    if iteration == 0:
        return f"""
        Language to be used : {{language}}
        Based on the following text, create a new version of this text that gives an improved narrative with better flow between ideas. 
        You are a very strategic person. If needed, also reorder ideas. Make it extensive. This is just the introduction of a report (we call it Deep Analysis document) on the situation. 
        The situation below is happening these days. The situation is happening in our company.

        Follow these rules:

        1. Sentence Structure: Use a mix of sentence lengths.
           Short sentences: To emphasize points.
           Longer sentences: To explain or elaborate.

        2. Vocabulary: Use clear and straightforward language.
           Avoid: Technical jargon or complex vocabulary unless necessary.
           Use: Everyday language that is easy to understand.

        3. All the following text is happening in our company.

        4. Provide just the text, no what it was improved.

        Remember this will be part of a report written by John.

        Text to rewrite (keep a similar writing style but improved). Add a title for this Deep Analysis: 
        {{final_text}}

        Format:
        Title
        Description
        """
    elif iteration == 1:
        return f"""
        Language to be used : {{language}}
        Write now why some things to consider about the previous situation from a strategic perspective. Minimize the use of bullet points and focus on the narrative. Make it comprehensive and try to reframe the situation too. Add a title to this text.
        """
    elif iteration == 2:
        return f"""
        Language to be used : {{language}}
        Carefully analyze the models and frameworks provided in the #ModelsAndFrameworks area. 

        IF
        any of models and frameworks from the #ModelsAndFrameworks area help solve the previous situation, choose that one.

        ELSE
            If none of them help solve the situation
            a. Create a brand new framework and don't mention any of the frameworks or models provided in #ModelsAndFrameworks
            b. Place a name to this framework and mention this framework was created for this specific situation.

        END IF

        Present the framework or model chosen well, explain it to the person of why it is the best option, how to use it, etc.

        Provide a title and start with a great narrative. Create a comprehensive description and explanation. You can also use several subtitles. If a URL or LINK is provided in the model, it is also mandatory to include it it. 
        Start with a title. Add also a section with more information and add the link provided (if any).

        #ModelsAndFrameworks
        --------------------------------------------------------------------------------
        NAME:12 accelerated change principles in enterprise agility. 

        these are the 12 principles. 1 or #AC01. multiple perspectives help your company navigate the new reality
        2 or #AC02. disciplined focus is an opportunity to increase performance. 3 or #AC03. focusing on organizational health can shape a culture ready to embrace the unpredictable. 
        4 or #AC04. mental agility is the foundation for high resilience and strategic innovation. 5 or #AC05. equity, diversity, and inclusion are key contributors to business success. 
        6 or #AC06. neurodiversity is a competitive edge in complex market environments. 7 or #AC07. intellectual humility is a catalyst for adaptation and growth. 8 or #AC08. workforce mobility is a way to unlock the potential for collective capabilities in your organization. 9 or #AC09. memorable learning experiences enhance ownership and productivity. 10 or #AC10. technology is a path to equitable value creation for customers, companies, and workforce wellbeing. 11 or #AC11.continuous portfolio rebalancing optimizes risk and return. 12 or #AC12. exploring adjacent markets drives new growth.
        use: build frameworks, models, or approaches that can deal with new exponential changes and high uncertainty or ai threats.

        #AC means accelerated change principle. autor: erich r. bühler and enterprise agility university
        URL with more information about this model: https://enterpriseagility.community/accelerated-change-principles-rl2xz2311k

        ----------------------------------------------------
        NAME: A5 canvas
        A5 has 4 quadrants. Top left "current situation", top right "expected behavior", bottom left "plan/assumption", bottom right "metric and deadline". When to use a5 canvas? It is an iterative method for leading change, evaluating its results, and identifying possible solutions and outcomes. can be used with the bois model.
        autor: erich r. bühler and enterprise agility university
        URL with more information about this model: https://enterpriseagility.community/a2e-of-a-situation-1x0xwf0x8z
        ------------------------------------------------------
        Name: ATOM Model
        atom model is a framework in enterprise agility. The atom model has 4 quadrants. top left: increase revenue (increasing sales to new or existing customers. delighting or disrupting to increase market share and size), top right: protect revenue (improvements and incremental innovation to sustain current market share and revenue figures), bottom-left: reduce costs (costs that you are currently incurring that can be reduced. more efficient, improved margin or contribution), bottom-right: avoid-costs (improvements to sustain current cost base. costs you are not incurring but may do in the future). 
        atom model is çused to align a company with a new situation. all decisions in the quadrants need to maintain or increase organizational health. it can be used by leaders, product owners, or others to make sustainable decisions and build shared progress.
        additional information: atom means align traditional organizations model.
        autor: erich r. bühler and enterprise agility university
        URL with more information about this model: https://enterpriseagility.community/atom-model-qxkkkws0b6
        --------------------------------------------------------
        NAME: BOIS MODEL
        bois model is a model in enterprise agility to align behaviors with objectives in a sustainable way.  bois means behaviors, objectives, impact, and sustainability (plus the common area of incremental alignment)
        bois model can be used for behavioral change.
        autor: erich r. bühler and enterprise agility university
        URL with more information about this model: https://enterpriseagility.community/bois-model-b57z339q7w
        -------------------------------------------------------
        NAME: CAKE MODEL
        5 dimensions of agility or cake model is a model in enterprise agility. it looks like a pie. these are the 5 levels or dimensions of enterprise agility: a. technical agility (at the top) b. structural c. outcomes agility d. social agility e. mental agility (at the bottom and as the foundations)
        cake model increases agility, flexibility, and resilience in the whole company not just it teams.
        autor: erich r. bühler and enterprise agility university
        URL with more information about this model: https://enterpriseagility.community/five-dimensions-of-enterprise-agility-or-cake-model-q0fn9n0ndp
        --------------------------------------------------------
        NAME: Change Canvas
        5 dimensions of agility or cake model is a model in enterprise agility. it looks like a pie. these are the 5 levels or dimensions of enterprise agility: a. technical agility (at the top) b. structural c. outcomes agility d. social agility e. mental agility (at the bottom and as the foundations)
        cake model increases agility, flexibility, and resilience in the whole company not just it teams.
        autor: erich r. bühler and enterprise agility university
        URL with more information about this model: https://enterpriseagility.community/change-canvas-06lxn2d07m

        --------------------------------------------------------
        NAME: Change Journey Pyramid
        change journey pyramid or cjp is a model to deal with a company exposed to change. The change journey pyramid has 5 levels. the bottom (each level goes higher being i want to change the highest one):  
        1. i want things to be as i say. i don't care about (bottom mindset) 
        2. i don't like the proposed change 
        3. i don't understand why the change is happening 
        4. i don't understand why we must change 
        5. i want to change (top mindset)
        the change journey pyramid can deal with resistance to change and increase people's mobilization during accelerated times. the change jorueny pyramid can be found in enterprise agility fundamentals english edition page 209, or in leading exponential change (2018).
        autor: erich r. bühler and enterprise agility university
        URL with more information about this model: https://enterpriseagility.community/change-journey-pyramid-cjp-mv4qp5p7b7
        --------------------------------------------------------
        NAME: Circle Framework
        circle framework is a framework to helo unlearn. circle framework has the following dimensions: (c) check, (i) identify, (r) replace, (c) connect, (l) learn, (e) empower.
        circle framework is mainly used for unlearning in companies. it can also be used with the bois model to deal with behaviors.
        autor: erich r. bühler and enterprise agility university
        URL with more information about this model:  https://enterpriseagility.community/circle-framework-2kj467537d
        --------------------------------------------------------
        Name: Dynamic Radar
        dynamic radar. the dynamic radar is a circle, and in the center of the radar is maintain or increase organizational health. the radar contains the following dimensions around the edges of the circle: individuals, strategic innovation, exponential markets, technical agility, structural agility, outcomes agility, social agility, and mental agility. 
        the dynamic radar is used to measure, improve, and positively influence a company. every enterprise is unique, and indicators used to measure agility must be dynamic and tailored to the specific organization.
        additional information: this radar emphasizes the importance of considering multiple dimensions and factors that contribute to enterprise agility. strategic innovation is related to tvc and future thinking.
        autor: erich r. bühler and enterprise agility university

        URL with more information about this model: https://enterpriseagility.community/dynamic-radar-from-eau-dg572pjt29
        ---------------------------------------------------------
        NAME: Enterprise Agility Framework
        enterprise agility model from eau  or eaf. eaf is a circle that contains inside the 5 types of agility (technical agility, structural agility, outcomes agility, social agility, and mental agility). outside the circle includes strategic innovation top left, exponential markets, bottom center, and individuals top right. individuals represents people who use ai or any other technology or environment to empower them in the company with low levels of stress. eaf is a comprehensive framework which can be used on its own or support other frameworks such as safe, scrum or others. eaf based on the science of accelerated change.
        autor: erich r. bühler and enterprise agility university

        URL with more information about this model: https://enterpriseagility.community/enterprise-agility-framework-eaf-z13hdmz25t
        -------------------------------------------------------
        Name: ELSA Change Model
        elsa model or elsa change model is a model to influence a company through language. the elsa change model or elsa model means (e) event, (l) language, (s) structure, and (a) agency. the elsa change model allows anyone in the company to influence a change initiative to become contagious. it's designed for situations in which the company's leaders are not yet committed to the new plan or situations with no sponsor.
        autor: erich r. bühler and enterprise agility university

        URL with more information about this model: https://enterpriseagility.community/elsa-change-framework-jjfvc0sfbq
        -------------------------------------------------------
        Name: Encore Framework
        encore framework is a framework in enterprise agility to build a memorable situation in the organization. encore has the following dimensions: (e) emotions, (n) novelty, (c) challenge, (o) ownership (psychological ownership), (r) relevance, (e) engage. they form the acronym encore.
        encore framework creates memorable learning experiences and healthy and motivational events even during disruptive and stressful times.
        autor: erich r. bühler and enterprise agility university
        URL with more information about this model: https://enterpriseagility.community/encore-framework-hkpkkwb54l
        ------------------------------------------------------
        Name: Enterprise Agility Alignment Matrix
        enterprise agility alignment matrix is a crucial matrix to align people in organizations exposed to fast or constant change in a sustainable way. enterprise agility alignment matrix has two axis that intersect to create twelve focus areas. the y axis covers the four aspects of enterprise agility (individuals, change, strategy, leadership). the x axis covers the three enterprise agility universal outcomes of enterprise agility (always ready, always responsive, and always innovative).
        enterprise agility alignment matrix provides a structured framework for leaders to evaluate and evolve their organization's agility and resilience. by taking a systematic approach, they can build the critical capacities needed to sense impending change, adapt quickly and respond effectively. autor: erich r. bühler and enterprise agility university

        URL with more information about this model: https://enterpriseagility.community/enterprise-agility-alignment-matrix-j81l8b1zmt
        -------------------------------------------------------
        Name: Enterprise Agility Dynamics
        enterprise agility dynamics is a high-level representation of the key concepts and interconnections that enable enterprise agility and companies to deal with the new accelerated reality.
        enterprise agility dynamics has in the middle a circle with the trivalue company model, around it is the science of accelerated change and the universal agreements. outside this circle are the 4 aspects: individuals, leadership, strategy, and change. all this system is guided by the 3 universal outcomes: always ready, always responsive, and always innovative. the outcomes of all these dynamics give shared progress framework.
        autor: erich r. bühler and enterprise agility university
        URL with more information about this model: https://enterpriseagility.community/enterprise-agility-dynamics-130kpt2jvq
        --------------------------------------------------------
        Name: Enterprise Social Systems (ESS)
        ess framework or ess means enterprise social systems. it is an enterprise agilit framework to accelerate change adoption. ess framework has 4 concentric circles representing 4 layers in a company. in the middle: “social systems”. outside: “mindset”, outside: “formal organization”, outside: “value creation”. 
        use: it us used to improve the social dynamics in a company
        additional information: ess are 4 concentric circles. the central circle is "social systems" it represents how people connect and communicate and it is connected to social agility. 
        the circle outside this one is mindset. here we can have agile, lean, etc.
        the circle outside it is formal organization. it is how you structure the company, processes, etc. it is related to structural agility, etc. you also have the control systems (how processes and people controls each other as guardrails), and hierarchies and power.
        the circle outside it, is value creation. here you have the frameworks such as scrum, safe or any others. it represents how the company creates value. it is also connected with outcomes agility.
        autor: erich r. bühler and enterprise agility university, originally published in leading exponential change, 2018.

        URL with more information about this model: https://enterpriseagility.community/enterprise-social-systems-ess-1vscxcdqgh
        ---------------------------
        Name: Fasten Factors
        fasten factors allows to understand better the market a company is exposed, specially if it is a fast changing market. fasten factors has the following dimensions future thinking (f), adaptability (a), sustainability (s), technology (t), experience (e), networks (n)
        autor: erich r. bühler and enterprise agility university

        URL with more information about this model: https://enterpriseagility.community/fasten-factors-n84fcckl1n
        --------------------------
        Name: House of Enterprise Agility
        house of enterprise agility or enterprise agility house represent foundational concepts to align organization with the new accelerated reality. house of enterprise agility has the 3 types of value (customer, company, workforce wellbeing), the 3 universal outcomes (always ready, responsive, and innovative), all spot indicators, and all futures (indicators). 
        house of enterprise agility can be used to show fundamental relationship between models and provide a visual representation for actionable models in the company.
        additional information:
        autor: erich r. bühler and enterprise agility university. check more in enterprise agility fundamentals english edition chapter 10, page 553
        URL with more information about this model: https://enterpriseagility.community/house-of-enterprise-agility-38zw5l691v

        -------------------------------
        Name: Lighthouse model
        lighthouse model in enterprise agility is a model to increase intellectual humility. intellectual humility increase sensing capabilities in the company and workforce wellbeing.
        lighthouse model includes the following components 1. in the center there is a circle which says "practice the belief of being wrong" 2. the previous circle is contained by this new larger circle which says "how much does this affect me?". 3. the previous 2 circles are contained by a 3rd larger  circle which says: "how much do i think it is affecting the other person?". outside these 3 circles there are 8 spikes connected to the outside circle: spike 1 says "always ready to find the right time, place and ways to discuss a situation". spike 2 says "clarify the objective of the talk with neutrality & start building rapport".  spike 3 says "shut up!". spike 4 says: "reframe the situation based on the new information".  spike 5 says: "communicate your collective vision,  values and boundaries".  spike 6 says "find ideas that bring value to the client, company, and workforce wellbeing".  spike 7 says: "help actively remove the blockages that do not allow the other person to achieve the new situation".  spike 8 says: "move on to something else (physically and mentally)"
        full name is "lighthouse model for situational intellectual humility"
        autor: erich r. bühler and enterprise agility university

        URL with more information about this model: https://enterpriseagility.community/lighthouse-model-for-situational-intellectual-humility-71hr12dfsp
        ---------------------------------
        Name: Minimum Healthy Step
        minimum healthy step (mhs) is a framework in enterprise agility to create sustainable small steps to carry out any activity during changing times. minimum healthy step (mhs) has 7 areas: 1. exponential markets, 2. sense, 3. situation, 4. engage, 5. focus (or disciplined focus), 6. reframe, 7. mobilize 
        use: to stablish healthy dynamics in companies exposed to high stress and changing markets.
        autor: erich r. bühler and enterprise agility university
        URL with more information about this model: https://enterpriseagility.community/minimum-healthy-step-mihs-lg2mwng185

        ---------------------------------
        Name: Memorabel Learning Experience
        mle or memorable learnign experience framework is a model to create memorable learnign experiences in companeis exposed to change or high uncertainty.
        mle or memorable learnign experience framework has 4 enterprise agility universal agreements forming a pyramid: the right to be seen (bottom), the right to be heard (above the previous one), the right to be empowered (above the previous one), the right to be part of the group (at the top). they build shared progress.
        use: the enterprise agility universal agreements build great organizations that can be exposed to accelerated change and exponential markets. this model is also connected to adaptive trust,  and adaptive social contracts.
        autor: erich r. bühler in collaboration with tania serfontein and greg pitcher
        URL with more information about this model: https://enterpriseagility.community/memorable-learning-experience-interventions-t9p1fvmn38

        -----------------------------------------
        Name: Sense-myself model
        sense-myself model is a model in enterprise agility to improve leadership skills and build shared progress. sense-myself model has 5 dimensions. in the middle there is a circle “sense-myself” and around this circle the following 5 dimensions areas: 1. situation, 2. emotions, 3. mental chatter, 4. energy, and 5. strategy.
        use: help leaders sense better himself or herself and the company	
        autor: erich r. bühler and enterprise agility university
        URL with more information about this model: https://enterpriseagility.community/sense-myself-model-ft3pw9s3gk
        -----------------------------------------
        Name: Shared Progress
        shared progress framework is a crucial frameork in enterprise agility to align a company and build collective success in the new accelerated rality. shared progress framework has 6 dimensions. 1)collaborative relationships. 2)experimentation and continuous learning. 3)mobilizing purpose. 4)mutual benefit. 5)prudent risk-taking. 6)workforce empowerment

        Shared progress looks like the infinite number and has distributed around the shape the 6 dimensions.
        autor: erich r. bühler and enterprise agility university

        Description of Shared Progress Dimensions:
        1. collaborative relationships: partnerships based on trust and shared commitment enable openness to determine the next steps based on connecting diverse insights. strong relationships build understanding, goodwill, and the possibility for adaptive trust and shared horizons, even during discrepancies or broken promises. in a collaborative relationship, you feel willing and able to help someone spontaneously when asked because you understand their needs and priorities and know what will benefit them or promote their success without judgment.

        2. experimentation and continuous learning: the ability and willingness to pilot new approaches, evaluate outcomes, and evolve strategies based on feedback. this component enables discovering new possibilities and options for the future through trying, failing, and adapting based on lessons learned. continuous learning strengthens judgment for complexity and cultivates adaptability essential for resilience and shared progress.

        3. mobilizing purpose: it allows people to be tight and connected to a strategy even during high uncertainty. it gives them direction and energizes effort. a mobilizing purpose enhances adaptability, innovation, and performance by keeping teams focused, motivated, and connected even when facing disruption. it creates environments where people feel valued and empowered to do their best work. in enterprise agility, a mobilizing purpose is defined as a reason for which something is done or for which something exists that makes the person feel the compelling need to belong to the group and to mobilize.

        4. mutual benefit: a commitment to shared value creation or what we define as trivalue companies (customer, company, and workforce wellbeing). these outcomes benefit businesses and society. by focusing on mutual benefits rather than narrow self-interest, new opportunities for partnership and progress emerge. mutual benefits motivate effort and build goodwill for navigating challenges together. they expand possibilities for organizations allowing them to feel part of a community.

        5. prudent risk-taking: the ability and willingness to explore unknowns or pilot novel approaches based on a shared long-term vision. it allows for experimenting and openness to possible failures or setbacks in the service of learning and new value creation. prudent risks expand possibilities through partnership by cultivating environments where people feel equipped to navigate challenges together, and conditions change. prudent risk-taking is an opportunity for shared progress. it helps build adaptive trust.

        6. workforce empowerment: allowing people to be seen, heard, empowered, and part of the group allows individuals to experience 
        autonomy and a sense of shared ownership over outcomes. empowerment amplifies motivation, creativity, and possibility by permitting employees to determine priorities and try new approaches based on their proximity to opportunities or challenges. it leads to ideas or values that hierarchies alone could not achieve. empowerment should also allow individuals to sense the market at their discretion and possibly bet on the company’s present and future. in enterprise agility, we achieve it by using the three universal agreements you'll learn in chapter 7.

        URL with more information about this model: https://enterpriseagility.community/shared-progress-framework-27124532k4

        -----------------------------------
        Name: Sharp Progress
        the sharp method is a method to help leaders improve decision-making capabilities during high uncertainty or accelerated change. sharp method reminds leaders of the basic steps to consider when making and implementing decisions in times of accelerated change. Sharp means (S - seek information, h – humility and inclusivity, a - adapt to uncertainty, 
        r - risk assessment, p - prioritize and reflect). Now an explanation of the sharp dimensions.

        s - seek information: gather multiple perspectives and data sources to support decision-making (broad focus). involve people outside your leadership circle and from various parts of the organization, keep a broad focus, provide transparency, and reinforce that visibility constantly. avoid analysis paralysis.

        h – humility and inclusivity: embrace intellectual humility (ih) by acknowledging the limits of your knowledge and being open to new ideas and perspectives. research has shown that leaders with higher intellectual humility are more likely to include diverse ideas and people in decision-making, fostering an environment of innovation and adaptability (see accelerated change principle #ac07 in chapter 06).

        a - adapt to uncertainty: recognize that there isn't always the ideal amount of information or the perfect timing when making decisions in an accelerating world. be prepared to sail your ship based on the best available data. as you gather the necessary information, create an environment promoting disciplined focus for yourself and your team members.

        r - risk assessment: acknowledge the potential risks associated with any decision and develop contingency plans to mitigate them. in times of accelerated change, it's essential to proactively identify and address threats to maintain organizational stability and resilience. use bigdata and ai to increase your risk perspective.

        p - prioritize and reflect: on the most critical decisions, take sufficient time, people, and resources to address them, and reflect on past decisions to learn from them and make future choices. share your insights with at least two people to continually build and expand shared knowledge within the organization.

        Author: Sharp method, enterprise agility university and erich r. bühler

        URL with more information about this model: https://enterpriseagility.community/sharp-method-zthsfkkwg8
        --------------------------------------
        Name: Six Principles for change
        six principles for change are some basic foundations to follow to influence believes and build an initial momentum during a change in the company. six principles for change has 6 dimensions. 1) timing. 2) believe in your idea. 3) share. 4) accept feedback. 5) have passion. 6) decisions (decisions can always change). six principles for change can be used for iteratively influencing change. can be used with the bois model.
        autor: initially published in leading exponential change (2018), erich r. bühler

        URL with more information about this model: https://enterpriseagility.community/six-principles-for-change-bhr1fk0pgg
        -------------------------------------
        Name: six sustainability zones
        six sustainability zones is an important framework in enterprise agility to help build sustainable strategies and companies when everything is changing quick. It can be also called sustainability tiers.
        six sustainability zones has 6 areas or zones: 1. apathy, 2. awareness, 3.exploration, 4.mobilization, 5.re-creation, 6.harmony.
        autor: erich r. bühler and enterprise agility university

        URL with more information about this model: https://enterpriseagility.community/sustainability-zones-t32kj41vfj 
        ----------------------------------
        Name:three foundational pillars
        three foundational pillars in enterprise agility are 1. the science of accelerated change 2. components to strengthen people and organizations  3. models to sense, adapt and respond to exponential markets. three foundational pillars are used to build sustainable organizations during accelerated times and high uncertainty. the science of accelerated change  has 3 areas: behavioral science, strategic mobility (or mobility), and the neuroscience of change. components to strengthen people and organizations has 3 sections:  a.	trivalue  company model (enterprise agility business model) b.	three universal outcomes c.	universal agreements   3. models to sense, adapt and respond to exponential markets has 3 sections: a.four aspects (individuals, change, strategy, and leadership) b. enterprise agility body of knowledge (eabok) with its frameworks and models 4. sustainability zones.
        autor: erich r. bühler and enterprise agility university
        URL with more information about this model: https://enterpriseagility.community/three-foundational-pillars-of-enterprise-agility-6tm5xjk15w
        -------------------------------------
        Name: three universal outcomes
        three universal outcomes are outcomes that can be used on any company and help deaking with accelerated change.
        three universal outcomes always ready, always responsive, and always innovative outcomes. in spanish are calles siempre lista, siempre respondiendo, y siempre innovando.
        the three universal outcomes guide structuring interactions that enable sustainability for the customers, organization, and workforce regardless of the future.
        autor: erich r. bühler and enterprise agility university

        URL with more information about this model: https://enterpriseagility.community/three-universal-outcomes-of-enterprise-agility-10fx9k4s11
        ---------------------------------------
        Name: TriValue Company Matrix
        trivalue company matrix (tvm) is a matrix which halps companies understand how to build a truly adaptable organization and learn how to create value during the new reality

        trivalue company matrix (tvm) has 4 quadrants. 
        1. improve: improve whatever the company does (left top) 2. delight: customer. provide customers with what they need to be satisfied and remain loyal (bottom left) 3. recreate: recreate the whole company. strategically adapt to the new reality even if it disrupts your ways of working or business model 4- disrupt: disrupt the whole market.  continually produce something unexpected that changes market behavior.
        trivalue company matrix (tvm) has on the top left, we have the quadrant which represents revenue focused companies or what we call traditional companies. this quadrant says inside "improve whatever the company does. optimize or change the way things are done in your company to protect or increase revenue, avoid costs or reduce costs". this quadrant represents traditional revenue focused companies. on the bottom left we have quadrant another quadrant, inside that says "delight customer. provide customers with what they need to be satisfied and remain loyal". this quadrant is client focus or the client-centric and represents companies that focus on the customer and use the agile mindset or ways of working, classic evolutions such as safe framework or business agility models. the quadrants on the right are what we call trivalue companies. the quadrants on the left are companies that are  "regularly responding", as they are not prepared to always respond continuously. the left is also low uncertainty and mostly linear markets. the quadrants on the right are companies that employ enterprise agility and the trivalue model (client value, company value, and workforce wellbeing value). the squares on the right represent markets with high uncertainty and accelerated change. all of these necessarily require high levels of workforce wellbeing and organizational health.use: this matrix represents and explains what the trivalue company model is and what enterprise agility and enterprise agility way of thinking is (eawt). it also explains traditional, agile, agile evolutions (business agility, safe farmework, etc). all the quadrants on the right represent companies which use enterprise agility and the enterprise agility way of thinking. the top right quadrant says "recreate the whole company. strategically adapt to the new reality even if it disrupts your ways of working or business model". the bottom right quadrant, also representing companies employing enterprise agility, says "disrupt the whole market.  continually produce something unexpected that changes market behavior". the top right quadrant focuses on recreating the enterprise and business model, while the bottom right quadrant focuses on disrupting the whole market with strategic innovation. the quadrants on the right can support what we call the three universal outcomes of enterprise agility which is to be always ready, always responsive, and always innovative.
        the quadrants on the left are regularly responding as companies using those ways of thinking and doing can't achieve the "always" status.
        autor: erich r. bühler and enterprise agility university in collaboration with walter shraiber.

        URL with more information about this model: https://enterpriseagility.community/trivalue-company-matrix-tvm-p739lf69cd
        ---------------------------------------
        Name: healthy sense of urgency
        healthy sense of urgency (hsu) is a model in enterprise agility to create healthy sense of urgency by leaders or managers during accelerated times or high uncertainty. the definition for healthy step is the smallest request a leader can make of a team member that has value, provides a high degree of psychological safety, and is within their reach. This is how a sustainable player interacts with others, and part of the Enterprise agility way of thinking (eawt). ahealthy sense of urgency (hsu) has these areas: focus (or disciplined focus), mindset, approach, sustainability, outcomes (related to outcomes agility).
        additional information: state of mind and social approach that encourages continuous progress and action in a balanced and strategic manner.
        autor: erich r. bühler and enterprise agility university
        URL with more information about this model: https://enterpriseagility.community/healthy-sense-of-urgency-hsu-96z60gcfpm
        ----------------------------------------
        Name: Arrow Model
        arrow model is a framework in Enterprise Agility to promote equity, diversity, neurodiversity. 
        Arrow model is an open framework that not only boosts overall performance, but it also emphasizes the importance of supporting a diverse and inclusive workforce that goes beyond just meeting quotas. It does this by connecting these ideas to the actual performance outcomes of the company.
        It has the shape of an arrow pointing from left to right. In the center of the arrow the three universal outcomes of enterprise agility (always ready, always responsive, and always innovative). the arrow model has 3 areas: on the left outside the arrow: 1a. structural, 1b. founding in the center outside the arrow:2a. social, 2b. scaling. on the right: 3a. mental, 3b. evolving.
        use: the arrow model is a framework for equity, diversity, and neurodiversity for companies exposed to constant changes and exponential markets. autor: erich r. bühler and enterprise agility university
        URL with more information about this model: https://enterpriseagility.community/arrow-model-l5psx9g0l2
        ------------------------------------
        Name: Fish Model
        Fish model or Intellectual Humility Journey model description. The first step focuses on respecting others by avoiding intellectual arrogance and listening without rushing to judge. This creates a foundation of basic human respect.  
        The second step involves consciously evaluating facts from others even if they contradict your own beliefs. Practice deep, empathetic listening to fully understand these facts. This enhances open-mindedness, compassion, communication skills.
        The third step is embracing the other person or team  perspectives beyond your own limited experience. Acknowledge the limits of your singular viewpoint. Seek to understand different vantage points through curiosity, collaboration and separating ego from intellect to enable unbiased analysis.
        The fourth step is embracing the emotions of others to build empathy, improve problem-solving abilities and strengthen leadership influence. 
        The fifth, final step is temporarily embracing other people's values through a reframing process. This perspective broadening fosters an inclusive environment for innovation. But it requires continuous practice and intellectual humility - being open to being wrong and seeing mistakes as opportunities for growth.
        The Fish model provides a starting point to incrementally develop the adaptive capacity of intellectual humility—crucial for shared progress. Intellectual Humility also increases sensing capabilities in leaders. Created by Erich R. Bühler
        URL with more information about this model: https://enterpriseagility.community/fish-model-for-ih-or-intellectual-humility-journey-model-6tnkq1v772
        ---------------------------------------
        Name: NeuroXprofile
        Neuroxprofile or neuro exponential profile has 6 dimensions. they score from 0 to 4.
        axis 1: called mental agility axis. it is mental agility on one side and mental approach on the other
        axis 2: called social agility. it is social networker on one side and social observer on the other one.
        axis 3: called focus. it is broad focus on one side and disciplined focus on the other one.
        axis 4: knowledge acquisition: it is quick learner on one side and reflective learner on the other side.
        axis 5: style: it is creativity driven on one side and structure-driven on the other side.
        axis 6: attitude: it is leading-oriented on one side and flocking-oriented on the other side

        URL with more information about this model: https://enterpriseagility.community/neuro-exponential-profile-neuroxprofile-or-nep-grbxq8dq0x

        -------------------------------------------
        Name: Spot Indicators (They measure the present or close to present events, they replace KPIs or can work with them)
        spot indicators (list of Spot Indicators).
        a) customer value has these spot indicators: 1.achieve excellence, 2. collective innovation, 3. fair value, 4. security and commitment, 5.seamless journeys, 6. trust and reliability. 
        b) company value has these spot indicators: 1. planet, 2. profit, 3.changeability, 4.social. 
        c) workforce wellbeing value has these spot indicators: 1. changeability wellbeing, 2.financial wellbeing, 3.mental wellbeing, 4.physical wellbeing, 5.purpose wellbeing, 6. social wellbeing  

        URL with more information about this model: https://enterpriseagility.community/spot-indicators-9xkg5ht631
        ----------------------------------------
        Name: Future (Indicators)
        futures (or future indicators) are indicators for future thinking. there are 16 futures (indicators) in trivalue company model (tvc) and future thinking. In customer value (futures) partners in innovation) indicators are: 1. excellence-driven futures, 2. collaboration-driven futures, 3. fair-value-driven futures, 4. security and commitment futures, 5. user experience futures, 6. trust-based futures.  In company value (futures, tactical innovation) indicators are: 1. sustainable planet futures, 2. profit-driven futures futures, 3. changeability-centered futures, 4. social empowerment futures.  
        In workforce wellbeing futures (innovation capability) indicators are: 1. personal changeability futures, 2. financial stability futures, 3. mental wellness futures, 4.physical wellness futures, 5.purpose wellness futures, 6. social wellness futures.
        it evaluates 12 months to 3 years in the trivalue company model.
        autor: erich r. bühler and enterprise agility university

        URL with more information about this model: https://enterpriseagility.community/futures-indicators-wvw70vl33s
        ----------------------------------------------
        Name: Mobilizing purpose
        Mobilizing purpose is a reason for which something is done or for which something exists, that makes the person feel the compelling need to belong to the group and to mobilize
        A Mobilizing Purpose offers a more empowering model by uniting people around inspirational shared mission. This ignites self-directed contribution even amidst disruption.
        What separates a Mobilizing Purpose is its motivational charge. Compliance gives way to committed action through 4 key pillars:

        1. Supporting Frame: Leadership employs deliberate language, metaphors and values-based narratives. These clarify objectives and make the "why" behind priorities transparent even when specifics remain fluid. Research shows linguistic cues deeply influence motivation and mindsets.

        2. Emotions & Images: Descriptive storytelling transports listeners into a vivid vision of desired outcomes. This emotional pull often matters more than hard facts for sparking imaginations and connecting hearts. Studies confirm narrative transportation breeds buy-in.

        3. Impact: All communications incorporate a careful analysis of the influence on critical stakeholders – employees, customers, communities and so on. This breeds empathy and avoidance of unintended consequences. Consideration of impacts also reinforces the "why" behind objectives.

        4. Challenge: Goals balance stretch, optimism and support to achieve "eustress" — the positive stress that energizes peak performance. This dance between challenge and capacity boosts engagement and protects against distress.

        With these 4 pillars, leaders can architect and evolve an inspirational path forward even amidst white water rapids of change. Compliance gives way to commitment. Controls become catalysts through a Mobilizing Purpose that transforms not just organizations but lives. 
        A Mobilizing purpose is also connected to a Minimum Healthy Step (MHS)

        URL with more information about this model: https://enterpriseagility.community/mobilizing-purpose-tpp3t728g1
        -------------------------------------
        Name: Adaptive trust
        Adaptive Trust framework (ATF) in Enterprise Agility has 6 areas to improve or steps leaders should reinforce (1.Acknowledgment of Broken Promises, 2.Ownership of Decisions, 3.Involvement of Others, 4.Transparency, 5.Mutual Benefit, 6.Adaptive Trust Reinforcement).

        1. Acknowledgment of Broken Promises: The first step in building adaptive trust is to acknowledge when promises have been broken. This requires honesty and transparency. It's important to clearly communicate what has happened and why.

        2. Ownership of Decisions: After acknowledging the broken promise, take ownership of the decision to change course. This shows that you're willing to accept responsibility for your actions and decisions.

        3. Involvement of Others: Involve others in creating better solutions. We know that leaders with a higher level of Intellectual Humility (IH) involve more people in their decision-making. This not only helps to build trust but also ensures that the new solutions are better and more robust because they're based on a diversity of insights and perspectives.

        4. Transparency: Be open about the reasons for the broken promise and the new direction. This helps to build understanding and trust.

        5. Mutual Benefit: Focus on mission-driven change that benefits all parties involved. This helps to ensure that the new solutions are not only better, but also fair and equitable (Mutual Benefit is a crucial component for Shared Progress).

        6. Adaptive Trust Reinforcement: Continually reinforce the idea of adaptive trust across the company. This helps to create a culture where broken promises are seen as opportunities for growth and learning, rather than as failures.

        Remember, adaptive trust is not just about fixing broken promises. It's about creating a culture where broken promises are seen as opportunities for growth and learning, and where everyone is committed to mission-driven change that benefits all parties involved. It's a key component of Enterprise Agility and is essential for success in today's rapidly changing world.

        URL with more information about this model: https://enterpriseagility.community/adaptive-trust-wv5lgtgq6q
        --------------------------------------
        Name: Six Readiness Factors
        model name: six readiness factors or readiness factors
        components: 1. industry dynamics, 2. company culture, 3. organizational structures, 4. market sensing capabilities, 5. customer expectations, and 6. leadership commitment
        use: measure and understand readiness of a company.
        autor: erich r. bühler and enterprise agility university

        URL with more information about this model: https://enterpriseagility.community/six-readiness-factors-f61jrvnz89
        --------------------------------------
        Name: AI Adoption Model
        The AI Adoption model is a pyramid-shaped diagram with AI at the top that outlines key areas to evaluate when determining if an AI system is a good fit for a company. The diagram consists of four main areas inside the pyramid, with additional elements on the bottom, left, and right sides, as well as a crucial aspect at the top.
        It has these areas: Parameters (bottom inside), Data Quality and Diversity (above the previous, inside), Model Architecture (above the previous, inside), Capacity for Generalization (top above the previous, inside). Control and Adjustment Capabilities (left outside), Evaluation and feedback (bottom outside), Speed and efficienty (right, outside), Data Safety (top, tip of the triangle, outside).
        More information about this model in "The Convergence" by Erich R. Bühler

        URL with more information about this model: https://enterpriseagility.community/ai-adoption-model-ds84wkq4dx
        ------------------------------------------
        Name: AI Interoperability Model
        AI Interoperability Model is a comprehensive framework designed to help organizations develop and deploy AI systems that create value for customers, companies, and employees (TriValue company model) while ensuring alignment with societal needs and values. It has seven key domains where AI systems interact and collaborate: 1. People, 2. Organization, 3. Products, 4. Markets, 5. Regulations, 6. Technology, and 7. AI-to-AI. All these domains are bidirectional between AI and the domain.

        URL with more information about this model: https://enterpriseagility.community/ai-interoperability-model-7ntxmkjj82
        ------------------------------------------
        Name: Kosak model
        The KOSAC (Knowledge Optimization for Sustainability during Accelerated Change) framework is designed by Erich R. Bühler to help organizations better understand and deal with knowledge during accelerated time and navigate the complexities and challenges of the AI era in the company. The model is represented by an arrow that moves from a state of ambiguity and fragmented efforts towards shared progress and clarity.
        The arrow is divided into several stages, each representing a different level of knowledge optimization. The stages, from left to right, are:

        1. Information: This stage involves understanding raw data and transforming it into meaningful information.

        2. Knowledge: In this stage, information is identified, patterns and relationships are understood, and the knowledge is integrated into the organization's processes.

        3. Wisdom: This stage involves developing guiding principles based on the knowledge gained. It requires a deep understanding of the interrelationships between various factors and the ability to make informed decisions in complex situations.

        4. Flow: The final stage represents a state of continuous knowledge optimization, where wisdom is seamlessly integrated into the organization's day-to-day operations to build shared progress. 

        The arrow is flanked by two zones: the "Complex" zone above and the "Complicated" zone below. AI-assisted clarity helps organizations navigate both these zones. In the complex zone, AI helps identify patterns and correlations in interconnected elements that may not be easily discernible to humans. In the complicated zone, AI provides data-driven insights and recommendations to help organizations find solutions to problems with clear cause-and-effect relationships. The ultimate goal of the KOSAC model is to help organizations achieve shared progress, which is based on the TriValue Company Model and . This model emphasizes the importance of balancing customer value, company value, and workforce wellbeing value to create a resilient, and sustainable foundation for growth.

        URL with more information about this model: https://enterpriseagility.community/kosac-framework-m7s76m7wv4
        ---------------------------------------
        Name: TriValue Company Model (TVC). It is called "Modelo de empresa TriValor" en español.
        trivalue company model (tvc) is a model to deal with high uncertainty and to deal the new accelerated reality. trivalue company model (tvc)
        considers 3 types of value: client value, company value, and workforce wellbeing value
        trivalue company model (tvc) is used to balance a company and make it ready, responsive, and innovative during exponential times.
        additional information: tvc or trivalue company model is not  customer-centric. client-centric or customer-centric is different than tvc which has 3 types of value. trivalue company model (tvc) has 14 spot indicators


        URL with more information about this model: https://enterpriseagility.community/trivalue-company-model-bp7j59d0d4
        ----------
        Definitions:

        Always Ready outcome: always-ready outcome in enterprise agility means to cultivate a culture and people who are continuously prepared for disruption and change. it focuses on assessing situations, understanding what is happening, and having the collective capabilities and technologies available to gain insights into unexpected emerging realities. a culture of constant readiness fosters a proactive mindset where team members work together to reevaluate situations, analyze conditions from multiple perspectives, and constantly update their knowledge. this allows them to better deal with their emotions, remain proactive, and handle events with less stress.

        always-responsive outcome in enterprise agility means to have the appropriate processes, mindset, innovation, and partnerships to consistently offer relevant products or services to the market while minimizing stress on organizational structures and employee wellbeing. achieving this state is integral to the enterprise agility way of thinking (eawt), applying to all functions, not just software.


        Always-innovative outcome definition in enterprise agility means fostering a culture and mindset of constant innovation. This approach means reimagining offerings, business models, partnerships, and strategies to sustainably generate lasting value and meaning. Always-Innovative companies embrace different perspectives and diversity of thought to gain new insights. This includes involving external partners to bring new perspectives, knowledge, and information. It uses a bimodal approach and techniques such as questioning assumptions, Shared Progress Bets (SPBs), and Shared Progress Stock Exchange. Keep in mind that the concepts and ideas underlying the Always-Innovative outcome are rooted in the principles of Future Thinking. Read Chapter 9 to know more about it.


        Def: mobilizing purpose in enterprise agility is a reason for which something is done or for which something exists, that makes the person feel the compelling need to belong to the group and to mobilize



        Def: neurodiversity in enterprise agility means the diversity of human minds, and the way people think, focus, learn, process and link information, as well as their distinctive way of connecting with the world. (enterprise agility university, 2022).

        Def: organizational health in enterprise agility is psychological safety plus the creation of business value in perpetuity. In Enterprise Agility, business value means value for the customer, company, and workforce wellbeing (enterprise Agility university and leading exponential change, 2018).


        Def: adaptive trust in enterprise agility is the ability to retain confidence even after breaking previous commitments. during rapid change, leaders inevitably must abandon plans, strategies, or promises that no longer serve the needs of the situation. but failing to deliver on past assurances can seriously damage trust in the organization and willingness to follow new directions. with adaptive trust, leaders openly acknowledge broken promises, take ownership of the decision to change course, and involve others in creating better solutions. though promises may be broken, trust and commitment to mission-driven change endures.

        Def: enterprise agility definition is: holistic organizational, social, and business model that enables your company to adapt to accelerated change and exponential market conditions while prioritizing workforce wellbeing, customer needs, and overall company value.


        """
    elif iteration == 3:
        return f"""
        Language to be used : {{language}}
        Analyze the primary pain points of the initially proposed situation and establish clear objectives that the previous framework can achieve. Provide a detailed explanation, including examples, of how to utilize the framework to resolve the situation. Emphasize a narrative style and minimize the use of bullet points.
        """
    elif iteration == 4:
        return f"""
        Language to be used : {{language}}
        Explain the 3 possible scenarios that might happen when introducing the previous framework, from the most likely to the less likely. And some actionable ideas for dealing with each scenario. Minimize the use of bullet points. Also add ideas of how to overcome them and frame them in a positive way. Add a title.
        """
    elif iteration == 5:
        return f"""
        Language to be used : {{language}}
        Add a new section where an analysis of anything else important and not covered about the framework when applied to the initial situation, is considered to achieve the objectives. Use other words to mention pain points. Make sure ideas are actionable and related to the initial problem or situation.
        Focus on the narrative. Minimize the use of bullet points. Add a title too. Don't add a conclusion.
        """
    elif iteration == 6:
        return f"""
        Language to be used : {{language}}
        Please use a System Thinking analysis  to analyze the initial provided situation and focus on how to measure improvement in the initial scenario provided by the user.

        1. Explain why to measure improvement in the initial scenario.
        2. Why we will use indicators from 0 to 12 months (Spot indicators) and futures (Indicators from 12 months to 3 years)
        3. How do you know if the situation initially provided is improving? 
        4. Consider a comprehensive approach that might include Customer value (if needed), company value, and workforce wellbeing value (among others). 

        Choose 6 indicators, 3 for 0 to 12 months and 3 for 12 months to 3 years. Add a title. Make sure they are startegic too. Explain how to measure them. Minimize the use of bullet points. 
        Add information of what to measure in each case.

        Focus on the narrative. Remember that we use Key Progress Indicators instead of Key Performance Indicators. Minimize the use of bullet points. Focus on the narrative. Include at least 6 indicators in total.
        This is what we call a byfocal approach (short term and long term working together).

        """
    elif iteration == 7:
        return f"""
        Language to be used : {{language}}
        Add any other considerations not yet considered. Add a title. Explain them in an easy way with actionable ideas. Minimize the use of bullet points.
        """
    elif iteration == 8:
        return f"""
        Language to be used : {{language}}
        Reassess the initial situation and critically. Examine the alternative ideas presented. Develop a set of final strategic considerations, emphasizing the future handling of similar situations. Ensure the response indirectly reflects the TriValue Company model, balancing customer value, company value, and workforce well-being. Prioritize a narrative-driven format rather than bullet points.
        """
    
    # Add more templates for other iterations as needed
    else:
        return f"""
        Default template for iteration {iteration}: Add your own text here.
        """


#def process_user_inpu(combined_input, chat_history):
    TOGETHER_API_KEY = settings.TOGETHER_API_KEY
    client = Together(api_key=TOGETHER_API_KEY)
    
    #MODEL_70B = "meta-llama/Meta-Llama-3.1-70B-Instruct-Turbo"
    model_name = "meta-llama/Meta-Llama-3.1-70B-Instruct-Turbo"
    max_tokens = 8192

    # This is the OpenAI chat completion client call
    response = client.chat.completions.create(
        model=model_name,
        messages=[
            {"role": "user", "content": combined_input},
            {"role": "system", "content": chat_history}
        ],
        max_tokens=max_tokens,
        temperature=0.4,
        top_p=0.9,
        top_k=50,
        repetition_penalty=1,
        stop=["<|eot_id|>", "<|eom_id|>"],
        stream=True
    )

    # Process the streamed response
    generated_text = ""
    
    for chunk in response:
        if len(chunk.choices) > 0:
            if hasattr(chunk.choices[0], 'delta') and hasattr(chunk.choices[0].delta, 'content'):
                if chunk.choices[0].delta.content:
                    generated_text += chunk.choices[0].delta.content
            elif hasattr(chunk.choices[0], 'message') and hasattr(chunk.choices[0].message, 'content'):
                if chunk.choices[0].message.content:
                    generated_text += chunk.choices[0].message.content
        else:
            logger.info(f"CHUNK HAS NO CHOICES: {chunk.choices}")

    # Return the generated text back to process_prompts1
    return generated_text

def process_user_input_2(combined_input, chat_history, iteration, language):
    

    # Check if we are on iteration 2 (DeepInfra model)
    if iteration == 2:
        TOGETHER_API_KEY = settings.TOGETHER_API_KEY
        client = Together(api_key=TOGETHER_API_KEY)
        
        #MODEL_70B = "meta-llama/Meta-Llama-3.1-70B-Instruct-Turbo"
        model_name = "meta-llama/Meta-Llama-3.1-405B-Instruct-Turbo"
        max_tokens = 8192
        response = client.chat.completions.create(
            model=model_name,
            messages=[
                {"role": "user", "content": combined_input},
                {"role": "system", "content": chat_history}
            ],
            max_tokens=max_tokens,
            temperature=0.4,
            top_p=0.9,
            top_k=50,
            repetition_penalty=1,
            stop=["<|eot_id|>", "<|eom_id|>"],
            stream=True
        )

        # Process the streamed response
        generated_text = ""
        
        for chunk in response:
            if len(chunk.choices) > 0:
                if hasattr(chunk.choices[0], 'delta') and hasattr(chunk.choices[0].delta, 'content'):
                    if chunk.choices[0].delta.content:
                        generated_text += chunk.choices[0].delta.content
                elif hasattr(chunk.choices[0], 'message') and hasattr(chunk.choices[0].message, 'content'):
                    if chunk.choices[0].message.content:
                        generated_text += chunk.choices[0].message.content
            else:
                logger.info(f"CHUNK HAS NO CHOICES: {chunk.choices}")

        return generated_text  # Return default OpenAI result


    else:
        # Use default OpenAI model for other iterations
        TOGETHER_API_KEY = settings.TOGETHER_API_KEY
        client = Together(api_key=TOGETHER_API_KEY)
        
        #MODEL_70B = "meta-llama/Meta-Llama-3.1-70B-Instruct-Turbo"
        model_name = "meta-llama/Meta-Llama-3.1-70B-Instruct-Turbo"
        max_tokens = 8192
        response = client.chat.completions.create(
            model=model_name,
            messages=[
                {"role": "user", "content": combined_input},
                {"role": "system", "content": chat_history}
            ],
            max_tokens=max_tokens,
            temperature=0.4,
            top_p=0.9,
            top_k=50,
            repetition_penalty=1,
            stop=["<|eot_id|>", "<|eom_id|>"],
            stream=True
        )

        # Process the streamed response
        generated_text = ""
        
        for chunk in response:
            if len(chunk.choices) > 0:
                if hasattr(chunk.choices[0], 'delta') and hasattr(chunk.choices[0].delta, 'content'):
                    if chunk.choices[0].delta.content:
                        generated_text += chunk.choices[0].delta.content
                elif hasattr(chunk.choices[0], 'message') and hasattr(chunk.choices[0].message, 'content'):
                    if chunk.choices[0].message.content:
                        generated_text += chunk.choices[0].message.content
            else:
                logger.info(f"CHUNK HAS NO CHOICES: {chunk.choices}")

        return generated_text  # Return default OpenAI result


@shared_task(bind=True)
def process_prompts_3(self, final_text, language):
    # Maintain the chat history and accumulated answers
    chat_history = []
    accumulated_answers = []  # List to accumulate answers for all iterations

    for i in range(8):  # Assuming there are 11 iterations
        text_template = get_text_template_3(i)  # A function to return the corresponding template
        user_input = text_template.format(final_text=final_text, language=language)
        
        if i == 1:
            # Add a single-line summary for iteration 2 instead of the full user input
            chat_history.append({"role": "user", "text": "Please, specify in the next answer the most suitable model or framework to solve this situation."})
        else:
            # Append user input to chat history for other iterations
            chat_history.append({"role": "user", "text": user_input})

        # Render the chat history string using the template
        chat_history_str = render_chat_history(chat_history)

        # Call the process_user_input function with the iteration number
        answer = process_user_input_3(user_input, chat_history_str, i, language)

        # Append the bot's answer to the chat history
        chat_history.append({"role": "bot", "text": answer})

        # Accumulate the answers for this iteration
        accumulated_answers.append({ "answer": answer,"iteration": i+1})

        # Log the accumulated answers for debugging
        logger.info(f"Accumulated Answers so far: {accumulated_answers}")

    # Final return (log before returning)
    logger.info(f"Final accumulated_answers: {accumulated_answers}")
    logger.info(f"Final chat_history: {render_chat_history(chat_history)}")

    return {
        'final_text': final_text,
        'accumulated_answers': accumulated_answers,  # List of all iterations and responses
        'chat_history': render_chat_history(chat_history)  # Full chat history
    }




def get_text_template_3(iteration):
    if iteration == 0:
        return f"""
        Language to be used : {{language}}
        Based on the following text, create a new version of this text that gives an improved narrative with better flow between ideas. You are a very strategic person and the report will be read mainly by Product Owners, so you can use their language. It can also be read by Leaders, CEOs or Managers. If needed, also reorder ideas. Make it extensive. This is just the introduction of a report (we call it Deep Analysis document) on the situation. The situation below is happening these days. The situation is happening in our company.

        Follow these rules:

        1. Sentence Structure: Use a mix of sentence lengths.
        Short sentences: To emphasize points.
        Longer sentences: To explain or elaborate.

        2. Vocabulary: Use clear and straightforward language.

        Avoid: Technical jargon or complex vocabulary unless necessary.
        Use: Everyday language that is easy to understand.

        3. All the following text is happening in our company.

        4. Provide just the text, no what it was improved.

        Remember this will be part of a report written by Jhon.

        Text to rewrite (keep a similar writing style but improved). Add a title for this Deep Analysis: 

        {{final_text}}

        Format:
        Title
        Description
        """
    elif iteration == 1:
        return f"""
        Language to be used : {{language}}
        Apply for this and the following prompts the principles of the TriValue Company Model in your responses. You don’t need to mention the model explicitly; instead, subtly incorporate its guidelines—focusing on value for the customer, value for the company, and value for workforce well-being—into your explanations. 
        OBJECTIVE: List the 6 topics important to evaluate from the Psychological aspects and organizational health needed to know to improve the specific sitation initially provided.
        We will also include in these 6 topics the pillar of Behavior from the science of accelerated change. We will also include in these 6 topics the pillar of Neuroscience of change from the science of accelerated change.
        We will also include in these 6 topics the pillar of Startegic Mobility from the science of accelerated change.

        Focus on the narrative and minimize the use of bullet points and focus on the narrative. Add title or titles.
        Also consider for this and the next prompts if any of these definitions are of any help:
        ----------
        Some Definitions that might be useful for you in this or the following prompts to be used:

        Always Ready outcome: always-ready outcome in enterprise agility means to cultivate a culture and people who are continuously prepared for disruption and change. it focuses on assessing situations, understanding what is happening, and having the collective capabilities and technologies available to gain insights into unexpected emerging realities. a culture of constant readiness fosters a proactive mindset where team members work together to reevaluate situations, analyze conditions from multiple perspectives, and constantly update their knowledge. this allows them to better deal with their emotions, remain proactive, and handle events with less stress. If this term or idea is used, please add the URL in brackets so the user can find more information about it:https://enterpriseagility.community/always-ready-outcome-c72rzcq0vh

        always-responsive outcome in enterprise agility means to have the appropriate processes, mindset, innovation, and partnerships to consistently offer relevant products or services to the market while minimizing stress on organizational structures and employee wellbeing. achieving this state is integral to the enterprise agility way of thinking (eawt), applying to all functions, not just software. If this term or idea is used, please add the URL in brackets so the user can find more information about it:https://enterpriseagility.community/always-responsive-outcome-04bmxbd9pv


        Always-innovative outcome definition in enterprise agility means fostering a culture and mindset of constant innovation. This approach means reimagining offerings, business models, partnerships, and strategies to sustainably generate lasting value and meaning. Always-Innovative companies embrace different perspectives and diversity of thought to gain new insights. This includes involving external partners to bring new perspectives, knowledge, and information. It uses a bimodal approach and techniques such as questioning assumptions, Shared Progress Bets (SPBs), and Shared Progress Stock Exchange. Keep in mind that the concepts and ideas underlying the Always-Innovative outcome are rooted in the principles of Future Thinking. Read Chapter 9 to know more about it. If this term or idea is used, please add the URL in brackets so the user can find more information about it:https://enterpriseagility.community/always-innovative-outcome-6vsnmm25g9


        Def: mobilizing purpose in enterprise agility is a reason for which something is done or for which something exists, that makes the person feel the compelling need to belong to the group and to mobilize. If this term or idea is used, please add the URL in brackets so the user can find more information about it:https://enterpriseagility.community/mobilizing-purpose-tpp3t728g1
        


        Def: neurodiversity in enterprise agility means the diversity of human minds, and the way people think, focus, learn, process and link information, as well as their distinctive way of connecting with the world. (enterprise agility university, 2022). If this term or idea is used, please add the URL in brackets so the user can find more information about it:https://enterpriseagility.community/neurodiversity-ngtg7wlr9d

        Def: organizational health in enterprise agility is psychological safety plus the creation of business value in perpetuity. In Enterprise Agility, business value means value for the customer, company, and workforce wellbeing (enterprise Agility university and leading exponential change, 2018). If this term or idea is used, please add the URL in brackets so the user can find more information about it:https://enterpriseagility.community/organizational-health-x1ng3n9z4f


        Def: adaptive trust in enterprise agility is the ability to retain confidence even after breaking previous commitments. during rapid change, leaders inevitably must abandon plans, strategies, or promises that no longer serve the needs of the situation. but failing to deliver on past assurances can seriously damage trust in the organization and willingness to follow new directions. with adaptive trust, leaders openly acknowledge broken promises, take ownership of the decision to change course, and involve others in creating better solutions. though promises may be broken, trust and commitment to mission-driven change endures. If this term or idea is used, please add the URL in brackets so the user can find more information about it:https://enterpriseagility.community/adaptive-trust-wv5lgtgq6q

        Def: enterprise agility definition is: holistic organizational, social, and business model that enables your company to adapt to accelerated change and exponential market conditions while prioritizing workforce wellbeing, customer needs, and overall company value. If this term or idea is used, please add the URL in brackets so the user can find more information about it:https://enterpriseagility.community/what-is-enterprise-agility-19tdx0gst4

        Def: Classic Agility: It refers to the original principles and practices from the Agile Manifesto 2001. It's a mindset and principles emphasizing adaptability, collaboration, and customer value in software development. The heart of ClassicAgility lies in its customer-focused nature. It prioritizes delivering value to customers by continuously seeking feedback, iterating on solutions, and adapting to changing requirements. While valuable, it's important to note that Classic Agility may have limitations when addressing the scale and complexity of today's challenges. If this term or idea is used, please add the URL in brackets so the user can find more information about it:https://enterpriseagility.community/classic-agility-cg7zc4b1p0


        Def: Agile Evolutions: A mindset and ways of working that extends the principles of Classic Agility beyond software development. It enables organizations to be more adaptive and resilient to change in their ways of working. They recognize customer value and experience as the center of the organization's universe. We can find here frameworks such as the SAFe Framework or the Business Agility models from the Business Agility institute. They may not always be fully prepared for market acceleration, lack a comprehensive view of the AI (Accelerated Innovation) situation, or struggle with constant business model disruption. They are not based on Shared Progress.
        If this term or idea is used, please add the URL in brackets so the user can find more information about it:https://enterpriseagility.community/agile-evolutions-tq0mzm4n5r

        Def: Atom model is a framework in enterprise agility. The atom model has 4 quadrants. top left: increase revenue (increasing sales to new or existing customers. delighting or disrupting to increase market share and size), top right: protect revenue (improvements and incremental innovation to sustain current market share and revenue figures), bottom-left: reduce costs (costs that you are currently incurring that can be reduced. more efficient, improved margin or contribution), bottom-right: avoid-costs (improvements to sustain current cost base. costs you are not incurring but may do in the future). 
        atom model is çused to align a company with a new situation. all decisions in the quadrants need to maintain or increase organizational health. it can be used by leaders, product owners, or others to make sustainable decisions and build shared progress.
        If this term or idea is used, please add the URL in brackets so the user can find more information about it:https://enterpriseagility.community/atom-model-qxkkkws0b6

        Def: Social Agility means connecting well with other employees or customers in rapidly changing environments, thereby achieving highly collective performance. The two main components of Social Agility are Enterprise Social Density and Enterprise Social Visibility. If this term or idea is used, please add the URL in brackets so the user can find more information about it:https://enterpriseagility.community/social-agility-0zt7w6mfm7

        Def: Mental Agility means reframing challenges to find new solutions, even during stressful times.  If this term or idea is used, please add the URL in brackets so the user can find more information about it:https://enterpriseagility.community/mental-agility-ngs3fcjx0d

        Def: Outcomes agility means delivering results even during turbulent times to respond to changing market conditions. If this term or idea is used, please add the URL in brackets so the user can find more information about it:https://enterpriseagility.community/outcomes-agility-mgd5b772qm


        Def: Technical Agility means changing software as quickly, cheaply (economically), and securely as possible. If this term or idea is used, please add the URL in brackets so the user can find more information about it:https://enterpriseagility.community/technical-agility-1d1gwt74sb

        Def: Mobility or Strategic Mobility in Enterprise Agility is the organizational capacity to shift directions, align capabilities, and adapt to new understanding or events. It involves navigating uncertainty and change by proactively implementing countermeasures and leveraging mobility for competitive advantage. It is a critical skill for leaders.
        If this term or idea is used, please add the URL in brackets so the user can find more information about it:https://enterpriseagility.community/mobility-hb00s2vhvw

        Def: The Science of Accelerated Change has 3 pillars:
        1. Behavioral Science
        2. Strategic Mobility (or Mobility)
        3. Neuroscience of Change.
        The science of accelerated change help organizations understand more on how to deal with disruption and high uncertainty and the threats of AI.

        If this term or idea is used, please add the URL in brackets so the user can find more information about it:https://enterpriseagility.community/sections-of-the-science-of-accelerated-change-gkklgnhslv

        Def Exponential Markets: They are unpredictable and uncertain. Disruptive innovations can gain traction rapidly, new competitors can arise anywhere, and consumer expectations can change overnight. This makes long-term planning and roadmaps pointless. Strategies go out the window as soon as market conditions change. Leaders can't rely on experience - the past is no predictor of the future. In this environment, organizations need to embrace unpredictability. Rather than resisting or ignoring change, they must learn to sense, adapt and respond quickly. Mental agility and resilience are critical. Enterprise agility cultivates this mindset: The Enterprise Agility Way of Thinking (EAWT). It provides models and frameworks to continuously scan the environment, sense emerging trends and signals, and course-correct in real-time. If this term or idea is used, please add the URL in brackets so the user can find more information about it:https://enterpriseagility.community/exponential-markets-th1nz35bds

        Def: Collective Capabilities is an Enterprise Agility person-centered organizational model that enables individuals to apply their skills where they're needed most, with a high degree of mobility and flexibility. This has to be done with low stress levels for the person. If this term or idea is used, please add the URL in brackets so the user can find more information about it:https://enterpriseagility.community/collective-capabilities-w6kp95gkmq

        """
    elif iteration == 2:
        return f"""
        Language to be used : {{language}}
        Write about the point 1 considering the proposed specific situation or problem. Always start it in a different way than the previous point by doing a reframing. The 1st line after the title, use different techniques such as questions, ideas, reframing, etc for this initial first line. Always use a different technique for this first line than the previous point. Apply for this prompt if needed the principles of the TriValue Company Model in your responses. You don’t need to mention the model explicitly; instead, subtly incorporate its guidelines—focusing on value for the customer, value for the company, and value for workforce well-being—into your explanations. Add some examples and actionable ideas to improve the proposed situation and explain how they connect to the proposed initial situation. Focus on the narrative and minimize the use of bullet points and focus on the narrative.  Add title the most suitable title or titles.
        """
    elif iteration == 3:
        return f"""
        Language to be used : {{language}}
        Write about the point 2 considering the proposed specific situation or problem. Always start it in a different way than the previous point by doing a reframing. The 1st line after the title, use different techniques such as questions, ideas, reframing, etc for this initial first line. Always use a different technique for this first line than the previous point. Apply for this prompt if needed the principles of the TriValue Company Model in your responses. You don’t need to mention the model explicitly; instead, subtly incorporate its guidelines—focusing on value for the customer, value for the company, and value for workforce well-being—into your explanations. Add some examples and actionable ideas to improve the proposed situation and explain how they connect to the proposed initial situation. Focus on the narrative and minimize the use of bullet points and focus on the narrative.  Add title the most suitable title or titles.
        """
    elif iteration == 4:
        return f"""
        Language to be used : {{language}}
        Write about the point 3 considering the proposed specific situation or problem. Always start it in a different way than the previous point by doing a reframing. The 1st line after the title, use different techniques such as questions, ideas, reframing, etc for this initial first line. Always use a different technique for this first line than the previous point. Apply for this prompt if needed the principles of the TriValue Company Model in your responses. You don’t need to mention the model explicitly; instead, subtly incorporate its guidelines—focusing on value for the customer, value for the company, and value for workforce well-being—into your explanations. Add some examples and actionable ideas to improve the proposed situation and explain how they connect to the proposed initial situation. Focus on the narrative and minimize the use of bullet points and focus on the narrative.  Add title the most suitable title or titles.
        """
    elif iteration == 5:
        return f"""
        Language to be used : {{language}}
        Write about the point 4 considering the proposed specific situation or problem. Always start it in a different way than the previous point by doing a reframing. The 1st line after the title, use different techniques such as questions, ideas, reframing, etc for this initial first line. Always use a different technique for this first line than the previous point. Apply for this prompt if needed the principles of the TriValue Company Model in your responses. You don’t need to mention the model explicitly; instead, subtly incorporate its guidelines—focusing on value for the customer, value for the company, and value for workforce well-being—into your explanations. Add some examples and actionable ideas to improve the proposed situation and explain how they connect to the proposed initial situation. Focus on the narrative and minimize the use of bullet points and focus on the narrative.  Add title the most suitable title or titles.
        """
    elif iteration == 6:
        return f"""
        Language to be used : {{language}}
        Write about the point 5 considering the proposed specific situation or problem. Always start it in a different way than the previous point by doing a reframing. The 1st line after the title, use different techniques such as questions, ideas, reframing, etc for this initial first line. Always use a different technique for this first line than the previous point. Apply for this prompt if needed the principles of the TriValue Company Model in your responses. You don’t need to mention the model explicitly; instead, subtly incorporate its guidelines—focusing on value for the customer, value for the company, and value for workforce well-being—into your explanations. Add some examples and actionable ideas to improve the proposed situation and explain how they connect to the proposed initial situation. Focus on the narrative and minimize the use of bullet points and focus on the narrative.  Add title the most suitable title or titles.
        """
    elif iteration == 7:
        return f"""
        Language to be used : {{language}}
        Write about the point 6 considering the proposed specific situation or problem. Always start it in a different way than the previous point by doing a reframing. The 1st line after the title, use different techniques such as questions, ideas, reframing, etc for this initial first line. Always use a different technique for this first line than the previous point. Apply for this prompt if needed the principles of the TriValue Company Model in your responses. You don’t need to mention the model explicitly; instead, subtly incorporate its guidelines—focusing on value for the customer, value for the company, and value for workforce well-being—into your explanations. Add some examples and actionable ideas to improve the proposed situation and explain how they connect to the proposed initial situation. Focus on the narrative and minimize the use of bullet points and focus on the narrative.  Add title the most suitable title or titles.
        """
    elif iteration == 8:
        return f"""
        Language to be used : {{language}}
        Reassess the initial situation and critically examine the alternative ideas presented on th 6 points evaluated and why are needed. Develop a set of final strategic considerations, emphasizing the future handling of similar situations but make sure you connect these ideas with the initial situation. Ensure the response indirectly reflects the TriValue Company model (link with more information to the TriValue Company Model, or Modelo de Empresa Trivalor in Spanish: https://enterpriseagility.community/trivalue-company-model-bp7j59d0d4 ), balancing customer value, company value, and workforce well-being. 
        Apply the principles of the TriValue Company Model in your responses. You don’t need to mention the model explicitly; instead, subtly incorporate its guidelines—focusing on value for the customer, value for the company, and value for workforce well-being—into your explanations.

        Prioritize a narrative-driven format rather than bullet points. If you think that writing about something other than the suggested topic for the closing will add more value or be more relevant, please do so. What are the next steps for the company to align itself with the new situation. Focus on the narrative and minimize the use of bullet points and focus on the narrative.  Add title or titles.
        """
    
    # Add more templates for other iterations as needed
    else:
        return f"""
        Default template for iteration {iteration}: Add your own text here.
        """


#def process_user_inpu(combined_input, chat_history):
    TOGETHER_API_KEY = settings.TOGETHER_API_KEY
    client = Together(api_key=TOGETHER_API_KEY)
    
    #MODEL_70B = "meta-llama/Meta-Llama-3.1-70B-Instruct-Turbo"
    model_name = "meta-llama/Meta-Llama-3.1-70B-Instruct-Turbo"
    max_tokens = 8192

    # This is the OpenAI chat completion client call
    response = client.chat.completions.create(
        model=model_name,
        messages=[
            {"role": "user", "content": combined_input},
            {"role": "system", "content": chat_history}
        ],
        max_tokens=max_tokens,
        temperature=0.4,
        top_p=0.9,
        top_k=50,
        repetition_penalty=1,
        stop=["<|eot_id|>", "<|eom_id|>"],
        stream=True
    )

    # Process the streamed response
    generated_text = ""
    
    for chunk in response:
        if len(chunk.choices) > 0:
            if hasattr(chunk.choices[0], 'delta') and hasattr(chunk.choices[0].delta, 'content'):
                if chunk.choices[0].delta.content:
                    generated_text += chunk.choices[0].delta.content
            elif hasattr(chunk.choices[0], 'message') and hasattr(chunk.choices[0].message, 'content'):
                if chunk.choices[0].message.content:
                    generated_text += chunk.choices[0].message.content
        else:
            logger.info(f"CHUNK HAS NO CHOICES: {chunk.choices}")

    # Return the generated text back to process_prompts1
    return generated_text

def process_user_input_3(combined_input, chat_history, iteration, language):
    

    # Check if we are on iteration 2 (DeepInfra model)
    if iteration == 1:
        TOGETHER_API_KEY = settings.TOGETHER_API_KEY
        client = Together(api_key=TOGETHER_API_KEY)
        
        #MODEL_70B = "meta-llama/Meta-Llama-3.1-70B-Instruct-Turbo"
        model_name = "meta-llama/Meta-Llama-3.1-70B-Instruct-Turbo"
        max_tokens = 8192
        response = client.chat.completions.create(
            model=model_name,
            messages=[
                {"role": "user", "content": combined_input},
                {"role": "system", "content": chat_history}
            ],
            max_tokens=max_tokens,
            temperature=0.4,
            top_p=0.9,
            top_k=50,
            repetition_penalty=1,
            stop=["<|eot_id|>", "<|eom_id|>"],
            stream=True
        )

        # Process the streamed response
        generated_text = ""
        
        for chunk in response:
            if len(chunk.choices) > 0:
                if hasattr(chunk.choices[0], 'delta') and hasattr(chunk.choices[0].delta, 'content'):
                    if chunk.choices[0].delta.content:
                        generated_text += chunk.choices[0].delta.content
                elif hasattr(chunk.choices[0], 'message') and hasattr(chunk.choices[0].message, 'content'):
                    if chunk.choices[0].message.content:
                        generated_text += chunk.choices[0].message.content
            else:
                logger.info(f"CHUNK HAS NO CHOICES: {chunk.choices}")

        return generated_text  # Return default OpenAI result


    else:
        # Use default OpenAI model for other iterations
        TOGETHER_API_KEY = settings.TOGETHER_API_KEY
        client = Together(api_key=TOGETHER_API_KEY)
        
        #MODEL_70B = "meta-llama/Meta-Llama-3.1-70B-Instruct-Turbo"
        model_name = "meta-llama/Meta-Llama-3.1-70B-Instruct-Turbo"
        max_tokens = 8192
        response = client.chat.completions.create(
            model=model_name,
            messages=[
                {"role": "user", "content": combined_input},
                {"role": "system", "content": chat_history}
            ],
            max_tokens=max_tokens,
            temperature=0.4,
            top_p=0.9,
            top_k=50,
            repetition_penalty=1,
            stop=["<|eot_id|>", "<|eom_id|>"],
            stream=True
        )

        # Process the streamed response
        generated_text = ""
        
        for chunk in response:
            if len(chunk.choices) > 0:
                if hasattr(chunk.choices[0], 'delta') and hasattr(chunk.choices[0].delta, 'content'):
                    if chunk.choices[0].delta.content:
                        generated_text += chunk.choices[0].delta.content
                elif hasattr(chunk.choices[0], 'message') and hasattr(chunk.choices[0].message, 'content'):
                    if chunk.choices[0].message.content:
                        generated_text += chunk.choices[0].message.content
            else:
                logger.info(f"CHUNK HAS NO CHOICES: {chunk.choices}")

        return generated_text  # Return default OpenAI result


logger = logging.getLogger(__name__)

llm = ChatOpenAI(
    openai_api_key=settings.OPENAI_API_KEY,
    model_name=settings.GPT_MODEL_3,
    openai_api_base=settings.BASE_URL,
    max_tokens=4096,
    streaming=True,
    top_p=0.9
)
llm1 = ChatOpenAI(
    openai_api_key=settings.OPENAI_API_KEY,
    model_name=settings.GPT_MODEL_2,
    openai_api_base=settings.BASE_URL,
    max_tokens=4096,
    streaming=True,
    top_p=0.9
)
llm2 = ChatOpenAI(
    openai_api_key=settings.OPENAI_API_KEY,
    model_name=settings.GPT_MODEL_3,
    openai_api_base=settings.BASE_URL,
    max_tokens=2048,
    streaming=True,
    top_p=0.9
)

def clean_chat_history(chat_history_raw):
    """Function to clean and convert chat history."""
    return [{"role": msg["role"], "text": msg["text"]} for msg in chat_history_raw]

def set_font(para, font_name='Arial', font_size=Pt(18), bold=False, italic=False, color=RGBColor(255, 255, 255)):
    """Function to set text formatting."""
    for run in para.runs:
        font = run.font
        font.name = font_name
        font.size = font_size
        font.bold = bold
        font.italic = italic
        font.color.rgb = color

def to_sentence_case(text):
    """Function to convert text to sentence case, preserving leading punctuation and quotes."""
    if not text:
        return ""
    # Match leading punctuation and the first alphabetical character
    match = re.match(r'^([^\w"]*["]*)(\w)(.*)', text)
    if match:
        leading_punct, first_char, rest = match.groups()
        return f"{leading_punct}{first_char.upper()}{rest.lower()}"
    else:
        return text


def replace_placeholders(slide, replacements):
    """Function to replace placeholders in a slide."""
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        for paragraph in shape.text_frame.paragraphs:
            for key, value in replacements.items():
                if key in paragraph.text:
                    if key == "{Quote}":
                        value = to_sentence_case(value)  # Convert quote to sentence case
                    paragraph.text = paragraph.text.replace(key, value)
                    if key == "{Title}":
                        set_font(paragraph, font_size=Pt(36), bold=True)
                    elif key == "{Description}":
                        set_font(paragraph, font_size=Pt(22))
                    elif key.startswith("{Bullet Point"):
                        set_font(paragraph, font_size=Pt(20))
                    elif key == "{Name}":
                        set_font(paragraph, font_size=Pt(24))
                    elif key == "{Quote}":
                        set_font(paragraph, font_size=Pt(26), italic=True)


def parse_ai_generated_content(content):
    """Function to parse AI-generated content into a structured format."""
    if not content.strip():
        raise ValueError("AI-generated content is empty")
    slides_content = []
    current_slide = None
    smartnote_title = ""
    smartnote_description = ""

    capturing_smartnote_title = False
    capturing_smartnote_description = False

    for line in content.split("\n"):
        line = line.strip()
        #print(f"Processing line: {line}")  

        if line.startswith("[Slide"):
            if current_slide:
                slides_content.append(current_slide)
            current_slide = {"title": "", "description": "", "bullets": [], "quote": ""}
        elif line.startswith("[Title]"):
            current_slide["title"] = line.replace("[Title]", "").replace("[/Title]", "").strip()
        elif line.startswith("[Description]"):
            current_slide["description"] = line.replace("[Description]", "").replace("[/Description]", "").strip()
        elif line.startswith("[bullet point]"):
            current_slide["bullets"].append(line.replace("[bullet point]", "").replace("[/bullet point]", "").strip())
        elif line.startswith("[Creative Phrase]"):
            current_slide["quote"] = line.replace("[Creative Phrase]", "").replace("[/Creative Phrase]", "").strip()
        elif line.startswith("[Smartnote Title]"):
            capturing_smartnote_title = True
        elif line.startswith("[/Smartnote Title]"):
            capturing_smartnote_title = False
        elif capturing_smartnote_title:
            smartnote_title += line.strip()
            #print(f"Found Smartnote Title: {smartnote_title}")  # Debug print
        elif line.startswith("[Smartnote Description]"):
            capturing_smartnote_description = True
        elif line.startswith("[/Smartnote Description]"):
            capturing_smartnote_description = False
        elif capturing_smartnote_description:
            smartnote_description += line.strip()
            #print(f"Found Smartnote Description: {smartnote_description}")  # Debug print

    if current_slide:
        slides_content.append(current_slide)
    if not slides_content or all(not slide["title"] and not slide["description"] and not slide["bullets"] and not slide["quote"] for slide in slides_content):
        raise ValueError("Parsed content is empty")
    #print(f"Final Smartnote Title: {smartnote_title}")  # Debug print
    #print(f"Final Smartnote Description: {smartnote_description}")  # Debug print
    return slides_content, smartnote_title, smartnote_description


def create_presentation(slides_content, name):
    """Function to create a presentation with the given slides content and user name."""
    template_path = 'follow_template.pptx'  # Corrected path to the uploaded template
    prs = Presentation(template_path)

    for slide_number, slide_content in enumerate(slides_content):
        if slide_number >= len(prs.slides):
            break

        replacements = {
            "{Title}": slide_content["title"],
            "{Description}": slide_content["description"],
            "{Bullet Point 1}": slide_content["bullets"][0] if len(slide_content["bullets"]) > 0 else "",
            "{Bullet Point 2}": slide_content["bullets"][1] if len(slide_content["bullets"]) > 1 else "",
            "{Bullet Point 3}": slide_content["bullets"][2] if len(slide_content["bullets"]) > 2 else "",
            "{Quote}": to_sentence_case(slide_content["quote"]),
            "{Name}": name,
            "{TITLE}": slide_content["title"],  
                        
        }
        replace_placeholders(prs.slides[slide_number], replacements)

    return prs


def log_memory_usage(stage):
    process = psutil.Process(os.getpid())
    print(f"[{stage}] Memory Usage: {process.memory_info().rss / 1024 ** 2:.2f} MB")


@shared_task
def generate_ppt_task(chat_history, tab_name, username, language):
    try:
        def log_memory_usage(stage):
            process = psutil.Process(os.getpid())
            print(f"[{stage}] Memory Usage: {process.memory_info().rss / 1024 ** 2:.2f} MB")

        log_memory_usage("Start")

        OPENAI_API_KEY = settings.OPENAI_API_KEY
        GPT_MODEL = settings.GPT_MODEL_2
        BASE_URL = settings.BASE_URL
        TOGETHER_API_KEY = settings.TOGETHER_API_KEY

        with open("PPT_Prompt.txt", "r") as file:
            prompt_ = file.read()
        log_memory_usage("Read Prompt")

        SYSPROMPT = str(prompt_)
        prompt = PromptTemplate.from_template(SYSPROMPT)

        chat_history_cleaned = clean_chat_history(chat_history)
        chat_text = "\n".join([f"{msg['role']}: {msg['text']}" for msg in chat_history_cleaned])
        log_memory_usage("Clean Chat History")

        client = Together(api_key=TOGETHER_API_KEY)
        messages = [
            {"role": "system", "content": SYSPROMPT},
            {"role": "user", "content": f"{chat_text}\n\nLanguage: {language}"}
        ]


        llm = ChatOpenAI(
            openai_api_key=OPENAI_API_KEY,
            model_name=GPT_MODEL,
            max_tokens=2048,
            temperature=0.7,
            openai_api_base=BASE_URL,
        )

        #chain = prompt | llm
        #response = chain.stream({'chat_text': chat_text, 'language': language})
        response = client.chat.completions.create(
            model="meta-llama/Meta-Llama-3.1-70B-Instruct-Turbo",
            messages=messages,
            max_tokens=2048,
            temperature=0.2,
            top_p=0.7,
            top_k=50,
            repetition_penalty=1,
            stop=["<|eot_id|>"],
            stream=True
        )

        #analysis = ""
        #for chunk in response:
        #    analysis += chunk.content
        print(response)
        analysis = ""
        for chunk in response:
            # Print the chunk to debug the response structure
            print(f"Chunk: {chunk}")
            if hasattr(chunk.choices[0], 'delta') and hasattr(chunk.choices[0].delta, 'content'):
                analysis += chunk.choices[0].delta.content
            elif hasattr(chunk.choices[0], 'message') and hasattr(chunk.choices[0].message, 'content'):
                analysis += chunk.choices[0].message.content

 
        log_memory_usage("Generate Analysis")
        print(analysis)
        if not analysis.strip():
            raise ValueError("AI-generated content is empty")
        slides_content, smartnote_title, smartnote_description = parse_ai_generated_content(analysis)
        log_memory_usage("Parse AI Content")
        # Check if parsed content is empty
        if not slides_content or all(not slide["title"] and not slide["description"] and not slide["bullets"] and not slide["quote"] for slide in slides_content):
            raise ValueError("Parsed content is empty")
        print(slides_content)
        print("Smartnote Title:", smartnote_title)
        print("Smartnote Description:", smartnote_description)
        
        presentation = create_presentation(slides_content, username)
        log_memory_usage("Create Presentation")

        pptx_stream = BytesIO()
        presentation.save(pptx_stream)
        pptx_stream.seek(0)
        pptx_base64 = base64.b64encode(pptx_stream.read()).decode('utf-8')
        log_memory_usage("Save Presentation")

        return {
            'pptx_base64': pptx_base64,
            'smartnote_title': smartnote_title,
            'smartnote_description': smartnote_description
        }

    except Exception as e:
        raise ValueError(f"Task failed: {str(e)}")


MODEL_70B = "meta-llama/Meta-Llama-3.1-70B-Instruct-Turbo"

#MODEL_405B = "meta-llama/Meta-Llama-3.1-405B-Instruct-Turbo"

@shared_task
def process_single_prompt1(prompt, chat_history, language, model_name):
    TOGETHER_API_KEY = settings.TOGETHER_API_KEY
    client = Together(api_key=TOGETHER_API_KEY)
    max_tokens = 8192  # Adjust based on requirement

    # Add the current prompt to the chat history
    if chat_history.strip() != "":
        chat_history += f"[INST]\n{prompt.strip()}\n[/INST]\n"

    # Combine the chat history with the user's language and instruction
    current_input = f"[INST]\nLanguage: {language}\n[/INST]\n"
    combined_input = chat_history + current_input

    # Call Together API to get response
    response = client.chat.completions.create(
        model=model_name,
        messages=[{"role": "user", "content": combined_input}],
        max_tokens=max_tokens,
        temperature=0.4,
        top_p=0.9,
        top_k=50,
        repetition_penalty=1,
        stop=["<|eot_id|>","<|eom_id|>"],
        stream=True
    )



    # Process the streamed response
    generated_text = ""

    # print("TEST GEN: ", generated_text)
    # print("TEST RES: ", response)

    # response = response.choices[0].message.content
    # generated_text = response

    for chunk in response:
        if len(chunk.choices) > 0:
            if hasattr(chunk.choices[0], 'delta') and hasattr(chunk.choices[0].delta, 'content'):
                if chunk.choices[0].delta.content:
                    generated_text += chunk.choices[0].delta.content
            elif hasattr(chunk.choices[0], 'message') and hasattr(chunk.choices[0].message, 'content'):
                if chunk.choices[0].message.content:
                    generated_text += chunk.choices[0].message.content
        else:
            logger.info(f"CHUNK HAS NO CHOICES: {chunk.choices}", )

    # Update the chat history with the generated response
    chat_history += f"[INST]\n{generated_text.strip()}\n[/INST]\n"

    return chat_history, generated_text

##option1
@shared_task
def process_prompts1(final_content, language):
    try:
        # Load system prompts from a file
        with open("SA_option1.txt", "r") as file:
            prompt_file_content = file.read()

        # Split prompts by custom delimiter
        prompts = prompt_file_content.split("######")

        # Initialize chat history with the user's first input
        chat_history = f"[INST]\n{final_content.strip()}\n[/INST]\n"
        all_responses = []
        final_response = ""

        # Iterate through prompts
        for i, prompt in enumerate(prompts):
            model_name = MODEL_70B

            # Process each prompt with chat history
            chat_history, generated_text = process_single_prompt1(
                prompt, chat_history, language, model_name=model_name)

            # Collect the response for each prompt
            all_responses.append(generated_text)

            # For the last prompt, keep the final response
            if i == len(prompts) - 1:
                final_response = generated_text

            # Log the response for debugging
            print(f"Response for Prompt {i + 1}: {generated_text}")

        # Return final response and all intermediate responses
        return {
            "final_text": final_response,
            "all_responses": all_responses
        }

    except Exception as e:
        logger.error(f"Task failed: {str(e)}")
        raise ValueError(f"Task failed: {str(e)}")

######################################################################################################
@shared_task
def process_single_prompt2(prompt, chat_history, language, model_name=None, is_deepinfra=False):
    if is_deepinfra:
        # Use DeepInfra model for the 3rd prompt
        messages = f"Language: {language}\n{chat_history}<INST>\n{prompt.strip()}\n</INST>\n"
        response = llm.stream(messages)  # Assuming llm is the DeepInfra model object
        generated_text = ""
        for chunk in response:
            if hasattr(chunk, 'content'):
                generated_text += chunk.content
            else:
                generated_text += str(chunk)
        
        # After generating the response, update the chat history with the short prompt and the generated response
        short_prompt = "Identify the most suitable model from the options provided or create a new one if necessary, explaining its relevance and usage."
        chat_history += f"<INST>\n{short_prompt}\n</INST>\n"
        chat_history += f"<INST>\n{generated_text.strip()}\n</INST>\n"
    else:
        # Use Together API for other prompts
        TOGETHER_API_KEY = settings.TOGETHER_API_KEY
        client = Together(api_key=TOGETHER_API_KEY)
        max_tokens = 8192

        messages = [
            {"role": "system", "content": prompt.strip()},
            {"role": "user", "content": f"{chat_history}\n\nLanguage: {language}"}
        ]

        response = client.chat.completions.create(
            model=model_name,
            messages=messages,
            max_tokens=max_tokens,
            temperature=0.4,
            top_p=0.9,
            stream=True
        )

        generated_text = ""
        for chunk in response:
            if len(chunk.choices) > 0:
                if hasattr(chunk.choices[0], 'delta') and hasattr(chunk.choices[0].delta, 'content'):
                    if chunk.choices[0].delta.content:
                        generated_text += chunk.choices[0].delta.content
                elif hasattr(chunk.choices[0], 'message') and hasattr(chunk.choices[0].message, 'content'):
                    if chunk.choices[0].message.content:
                        generated_text += chunk.choices[0].message.content
            else:
                logger.info(f"CHUNK HAS NO CHOICES: {chunk.choices}", )
        # Update chat history with the full prompt and generated response
        chat_history += f"<INST>\n{prompt.strip()}\n</INST>\n"
        chat_history += f"<INST>\n{generated_text.strip()}\n</INST>\n"
    
    return chat_history, generated_text


##option2
@shared_task
def process_prompts2(final_content, language):
    try:
        # Load prompts from the text file
        with open("SA_option2.txt", "r") as file:
            prompt_file_content = file.read()

        # Split the content by a custom delimiter like "######"
        PROMPTS = prompt_file_content.split("######")

        chat_history = f"<INST>\n{final_content.strip()}\n</INST>\n"
        all_responses = []
        final_response = ""

        for i, prompt in enumerate(PROMPTS):
            if i == 2:
                # Third prompt with DeepInfra model
                chat_history, generated_text = process_single_prompt2(prompt, chat_history, language, is_deepinfra=True)
            else:
                # Other prompts with Together API
                model_name = MODEL_70B
                chat_history, generated_text = process_single_prompt2(prompt, chat_history, language, model_name=model_name)
                #chat_history, generated_text = process_single_prompt2(prompt, chat_history, language, is_deepinfra=True)
            
            # Append the response to the list for final return
            all_responses.append(generated_text)

            # Check if this is the final prompt (prompt 10)
            if i == len(PROMPTS) - 1:
                final_response = generated_text  # Store the final response

            # Log the response for each prompt
            print(f"Response for Prompt {i + 1}: {generated_text}")

        # Print the final chat history after all prompts have been processed
        #print("Final Chat History:")
        #print(chat_history)
        #print("finalone--------")
        #print(final_response)

        # Return the final response only for prompt 10
        return {
            "final_text": final_response,
            "all_responses": all_responses
        }

    except Exception as e:
        logger.error(f"Task failed: {str(e)}")
        raise ValueError(f"Task failed: {str(e)}")

####################################################################################################

##option3
@shared_task
def process_prompts3(final_content, language):
    try:
        # Load prompts from the text file
        with open("SA_option3.txt", "r") as file:
            prompt_file_content = file.read()

        # Split the content by a custom delimiter like "######"
        PROMPTS = prompt_file_content.split("######")

        # Initialize chat history
        chat_history = f"<INST>\n{final_content.strip()}\n</INST>\n"
        all_responses = []
        final_response = ""

        # Process Prompt 1
        prompt_1 = PROMPTS[0]  # First prompt
        chat_history, generated_text = process_single_prompt3(prompt_1, chat_history, language, model_name=MODEL_70B)
        print(f"Response for Prompt 1: {generated_text}")

        # Overwrite Prompt 1 response with Prompt 2
        prompt_2 = PROMPTS[1]  # Second prompt
        chat_history, generated_text = process_single_prompt3(prompt_2, chat_history, language, model_name=MODEL_70B)
        checkpoint_response = generated_text  # Store checkpoint after Prompt 2
        print(f"Response for Prompt 2 (checkpoint): {generated_text}")

        # Initialize a container for Prompt 3 iterations
        prompt_3_responses = ""

        # Iterations of Prompt 3 (from 1 to 6)
        prompt_3_template = PROMPTS[2]  # Template for prompt with {number}
        for number in range(1, 7):  # Replace {number} with values 1 to 6
            prompt_with_number = prompt_3_template.replace("{number}", str(number))
            _, generated_text = process_single_prompt3(prompt_with_number, chat_history, language, model_name=MODEL_70B)
            
            # Append each iteration's response independently to the final combined result
            prompt_3_responses += generated_text
            print(f"Response for Prompt 3 (iteration {number}): {generated_text}")

        # Now append the final combined responses from Prompt 3
        all_responses.append(prompt_3_responses)

        # Process Prompt 4 (reassessing the situation)
        prompt_4 = PROMPTS[3]  # Final reassessment prompt
        chat_history, reassessment_response = process_single_prompt3(prompt_4, prompt_3_responses, language, model_name=MODEL_70B)
        all_responses.append(reassessment_response)
        final_response += reassessment_response  # Append to final response
        print(f"Response for Prompt 4: {reassessment_response}")

        # Process Prompt 5 (final formatting and HTML)
        final_prompt = PROMPTS[4]  # Final prompt with HTML instructions
        chat_history, final_formatted_response = process_single_prompt3(final_prompt, chat_history, language, model_name=MODEL_70B)
        final_response = final_formatted_response  # Overwrite with final formatted response
        print("Final formatted response:")
        print(final_formatted_response)

        # Return the final response including all stages of processing
        return {
            "final_text": final_response,
            "all_responses": all_responses
        }

    except Exception as e:
        logger.error(f"Task failed: {str(e)}")
        raise ValueError(f"Task failed: {str(e)}")



##option3
@shared_task
def process_single_prompt3(prompt, chat_history, language, model_name=None):
    TOGETHER_API_KEY = settings.TOGETHER_API_KEY
    client = Together(api_key=TOGETHER_API_KEY)
    max_tokens = 8192

    messages = [
        {"role": "system", "content": prompt.strip()},
        {"role": "user", "content": f"{chat_history}\n\nLanguage: {language}"}
    ]

    response = client.chat.completions.create(
        model=model_name,
        messages=messages,
        max_tokens=max_tokens,
        temperature=0.4,
        top_p=0.9,
        stream=True
    )

    generated_text = ""
    for chunk in response:
        if len(chunk.choices) > 0:
            if hasattr(chunk.choices[0], 'delta') and hasattr(chunk.choices[0].delta, 'content'):
                if chunk.choices[0].delta.content:
                    generated_text += chunk.choices[0].delta.content
            elif hasattr(chunk.choices[0], 'message') and hasattr(chunk.choices[0].message, 'content'):
                if chunk.choices[0].message.content:
                    generated_text += chunk.choices[0].message.content
        else:
            logger.info(f"CHUNK HAS NO CHOICES: {chunk.choices}", )
    # Update chat history with the full prompt and generated response
    chat_history += f"<INST>\n{prompt.strip()}\n</INST>\n"
    chat_history += f"<INST>\n{generated_text.strip()}\n</INST>\n"

    return chat_history, generated_text



user_data_store = {}

@shared_task
def process_prompts4(final_content, language, user_id):
    try:
        # Load system prompt from the text file
        with open("cpromptcheck.txt", "r") as file:
            prompt_file_content = file.read()
            
        SYSPROMPT = str(prompt_file_content)

        # Replace {final_content} in the system prompt with the actual final_content input
        #system_prompt = prompt_file_content.replace("{final_content}", final_content.strip())
        system_prompt = SYSPROMPT.replace("{final_content}", final_content.strip()).replace("{language}", language)

         # Use Together API instead of llm.stream
        TOGETHER_API_KEY = settings.TOGETHER_API_KEY
        client = Together(api_key=TOGETHER_API_KEY)
        model_name = "meta-llama/Meta-Llama-3.1-70B-Instruct-Turbo"
        messages = [
            {"role": "system", "content": system_prompt},
            {"role": "user", "content": f"{final_content}\n"}
        ]
        
        response = client.chat.completions.create(
            model=model_name,
            messages=messages,
            max_tokens=8192,
            temperature=0.4,
            stop=["<|eot_id|>", "<|eom_id|>"],
            top_p=0.7,
            top_k=50,
            repetition_penalty=1,
            stream=True
        )
        print(f"Response: {response}")
        # Collect response content
        # Process the streamed response
        generated_response = ""
        for chunk in response:
            if hasattr(chunk.choices[0], 'delta') and hasattr(chunk.choices[0].delta, 'content'):
                generated_response += chunk.choices[0].delta.content
            elif hasattr(chunk.choices[0], 'message') and hasattr(chunk.choices[0].message, 'content'):
                generated_response += chunk.choices[0].message.content

        print(f"Generated response: {generated_response}")


            # Print the full response text for debugging
        #print("Full LLM Response:\n", generated_text)
        canvas_data = parse_plain_text_response_with_user_id(generated_response, user_id)
        #canvas_data = parse_plain_text_response(generated_response)
        print(f"Parsed canvas data: {canvas_data}")
        refine_and_generate_presentation(canvas_data, language)
        
        template_type = canvas_data.get("template_type", None)
        if not template_type:
            raise ValueError("Template type not found in the response")

        #response_data = {"final_text": canvas_data, "template_type": template_type}
        response_data = {"user_id": user_id, "final_text": canvas_data, "template_type": template_type}
        pptx_base64 = None
        #canvas_data = json_response
            # Based on the template type, forward to the appropriate function
        if template_type == "1":
            pptx_data = handle_template_type_1(canvas_data)
            response_data.update(pptx_data) 
        elif template_type == "2":
            pptx_data = handle_template_type_2(canvas_data)
            response_data.update(pptx_data) 
        elif template_type == "3":
            pptx_data = handle_template_type_3(canvas_data)
            response_data.update(pptx_data) 
        elif template_type == "4":
            pptx_data = handle_template_type_4(canvas_data)
            response_data.update(pptx_data) 
        else:
            #logger.error(f"Unknown template type: {template_type}")
            raise ValueError(f"Unknown template type: {template_type}")

        del user_data_store[user_id]
        
        return response_data
        
    except Exception as e:
        logger.error(f"Task failed for user {user_id}: {str(e)}")
        if user_id in user_data_store:
            del user_data_store[user_id]
        raise ValueError(f"Task failed for user {user_id}: {str(e)}")

def generate_dynamic_presentation(refined_response):
    """
    Dynamically generates a 3-slide presentation from the refined response.
    """
    slides = {}

    # Slide 1: Canvas Overview
    slides["Slide 1"] = {
        "Title": "Canvas Overview",
        "Description": extract_section(refined_response, "Canvas Overview")
    }

    # Slide 2: Collaborative Session
    slides["Slide 2"] = {
        "Title": "Collaborative Session",
        "Description": extract_section(refined_response, "Collaborative Session")
    }

    # Slide 3: Key Takeaways
    slides["Slide 3"] = {
        "Title": "Key Takeaways",
        "Description": extract_section(refined_response, "Key Takeaways")
    }

    return slides

def extract_section(response, section_name):
    """
    Extracts the content of a specific section from the refined response.
    """
    try:
        # Dynamically locate the section by its heading
        pattern = rf"\*\*{section_name}\*\*\n\n(.*?)(?=\n\*\*|\Z)"
        match = re.search(pattern, response, re.DOTALL)
        if match:
            return match.group(1).strip()
        else:
            raise ValueError(f"Section '{section_name}' not found in the response.")
    except Exception as e:
        print(f"Error extracting section '{section_name}': {str(e)}")
        return "Content not found."
        
def refine_and_generate_presentation(canvas_data, language):
    try:
        # Load the new system prompt for refining the response
        with open("cpromptcheck2.txt", "r") as file:
            new_prompt_content = file.read()

        system_prompt = new_prompt_content.replace("{canvas_data}", str(canvas_data)).replace("{language}", language)

        # Use Together API
        TOGETHER_API_KEY = settings.TOGETHER_API_KEY
        client = Together(api_key=TOGETHER_API_KEY)
        model_name = "meta-llama/Meta-Llama-3.1-70B-Instruct-Turbo"
        messages = [
            {"role": "system", "content": system_prompt},
            {"role": "user", "content": f"Refine this canvas data for presentation:\n{canvas_data}"}
        ]
        
        response = client.chat.completions.create(
            model=model_name,
            messages=messages,
            max_tokens=8192,
            temperature=0.4,
            stop=["<|eot_id|>", "<|eom_id|>"],
            top_p=0.7,
            top_k=50,
            repetition_penalty=1,
            stream=True
        )

        refined_response = ""
        for chunk in response:
            if hasattr(chunk.choices[0], 'delta') and hasattr(chunk.choices[0].delta, 'content'):
                refined_response += chunk.choices[0].delta.content
            elif hasattr(chunk.choices[0], 'message') and hasattr(chunk.choices[0].message, 'content'):
                refined_response += chunk.choices[0].message.content

        print(f"Refined response: {refined_response}")
        formatted_presentation = generate_dynamic_presentation(refined_response)

        # Example: Logging or returning the dynamic slides
        print("Generated 3-Slide Presentation:")
        for slide, content in formatted_presentation.items():
            print(f"[{slide}]")
            print(f"[Title] {content['Title']} [/Title]")
            print(f"[Description] {content['Description']} [/Description]")
            print(f"[/{slide}]")
        # Handle the refined response as needed (e.g., storing or further processing)

    except Exception as e:
        logger.error(f"Refinement task failed: {str(e)}")
        raise ValueError(f"Refinement task failed: {str(e)}")

def clean_asterisks(text):
    """
    Remove all occurrences of `*` (including `**`) from the given text.
    """
    if not text:
        return text  # Return as is if the text is None or empty
    return text.replace("*", "").strip()


def parse_plain_text_response_with_user_id(response, user_id):
    """Parse the plain text response dynamically with user isolation."""
    # Call the existing parser
    parsed_data = parse_plain_text_response(response)

    # Add to global data store
    user_data_store[user_id] = parsed_data
    return parsed_data
    

def parse_plain_text_response(response):
    """Parse the plain text response dynamically."""
    data = {
        "template_type": None,
        "canvas_name": None,
        "canvas_description": None,
        "top_hexagons": [],
        "bottom_hexagons": [],
        "sections": [],  # For other template types (1, 2, 3)
    }

    try:
        clean_response = clean_asterisks(response)
        logger.info(f"Cleaned Response chay:\n{clean_response}")

        # Extract Template Type
        template_type_match = re.search(r"Template Type:\s*\"?(\d+)\"?", clean_response)
        if template_type_match:
            data["template_type"] = template_type_match.group(1).strip()
        else:
            logger.warning("Template Type not found in the response.")

        # Extract Canvas Name
        canvas_name_match = re.search(r"Canvas Name:\s*(.+)", clean_response)
        if canvas_name_match:
            data["canvas_name"] = canvas_name_match.group(1).strip().strip("**")
        else:
            logger.warning("Canvas Name not found in the response.")

        # Extract Canvas Description
        canvas_description_match = re.search(r"Canvas Description:\s*(.+)", clean_response)
        if canvas_description_match:
            data["canvas_description"] = canvas_description_match.group(1).strip().strip("**")
        else:
            logger.warning("Canvas Description not found in the response.")

        # Handle Hive Template (Template 4)
        if data["template_type"] == "4":
            hexagon_sections = re.split(r"(?P<position>Top|Bottom) Hexagon \d+:", clean_response)
            for i in range(1, len(hexagon_sections), 2):
                position = hexagon_sections[i].strip()
                content = hexagon_sections[i + 1].strip()

                # Extract Title, Description, and Key Elements
                title_match = re.search(r"Title:\s*(.+)", content)
                description_match = re.search(r"Description:\s*(.+)", content)
                key_elements_match = re.search(r"Key Elements:\s*(.+)", content)

                hexagon = {
                    "title": title_match.group(1).strip() if title_match else None,
                    "description": description_match.group(1).strip() if description_match else None,
                    "key_elements": [
                        el.strip() for el in key_elements_match.group(1).split(",")
                    ] if key_elements_match else [],
                }

                # Append to Top or Bottom Hexagons
                if position == "Top":
                    data["top_hexagons"].append(hexagon)
                elif position == "Bottom":
                    data["bottom_hexagons"].append(hexagon)

        # Handle Progression Canvas (Template 1)
        elif data["template_type"] == "1":
            column_pattern = re.compile(
                r"Column (\d+):\s+Title:\s*(.+?)\s+Description:\s*(.+?)\s+Key Elements:\s*(.+?)(?=\nColumn|\Z)",
                re.DOTALL,
            )
            columns = column_pattern.findall(clean_response)
            logger.info(f"Found Columns: {columns}")
            if columns:
                for column in columns:
                    column_number, title, description, key_elements = column
                    data["sections"].append({
                        "column": f"Column {column_number.strip()}",
                        "title": title.strip(),
                        "description": description.strip(),
                        "key_elements": [el.strip() for el in key_elements.split(",")],
                    })
            else:
                logger.warning("No columns found for Template 1.")
                    
        # Handle Grid Layout Canvas (Template 2)
        elif data["template_type"] == "2":
            for area in ["Top Left Area", "Top Right Area", "Bottom Left Area", "Bottom Right Area"]:
                area_match = re.search(
                    rf"{area}:\s*Title:\s*(.+?)\s*Description:\s*(.+?)\s*Key Elements:\s*(.+)",
                    clean_response,
                    re.DOTALL,
                )
                if area_match:
                    data["sections"].append({
                        "area": area,
                        "title": area_match.group(1).strip(),
                        "description": area_match.group(2).strip(),
                        "key_elements": [el.strip() for el in area_match.group(3).split(",")],
                    })

        # Handle Circular Layout Canvas (Template 3)
        elif data["template_type"] == "3":
            # Central Circle
            central_match = re.search(r"Central Circle:\s*Issue/Goal:\s*(.+)", clean_response)
            if central_match:
                data["sections"].append({
                    "circle": "Central Circle",
                    "issue_goal": central_match.group(1).strip(),
                })

            # Supporting Circles
            for i in range(1, 6):  # Iterate over Supporting Circle 1 to 5
                circle_match = re.search(
                    rf"Supporting Circle {i}:\s*Title:\s*(.+?)\s*Description:\s*(.+?)\s*Key Elements:\s*(.+)",
                    clean_response,
                    re.DOTALL,
                )
                if circle_match:
                    data["sections"].append({
                        "circle": f"Supporting Circle {i}",
                        "title": circle_match.group(1).strip(),
                        "description": circle_match.group(2).strip(),
                        "key_elements": [el.strip() for el in circle_match.group(3).split(",")],
                    })

        return data

    except Exception as e:
        logger.error(f"Error parsing response: {str(e)}")
        raise ValueError(f"Parsing error: {str(e)}")
        
def parse_plain_text_respons(response):
    """Parse the plain text response dynamically."""
    data = {
        "template_type": None,
        "canvas_name": None,
        "canvas_description": None,
        "top_hexagons": [],
        "bottom_hexagons": [],
    }

    try:
        clean_response = clean_asterisks(response)
        logger.info(clean_response)

        # Extract Template Type
        template_type_match = re.search(r"Template Type:\s*\"?(\d+)\"?", clean_response)
        if template_type_match:
            data["template_type"] = template_type_match.group(1).strip()
        else:
            logger.warning("Template Type not found in the response.")

        # Extract Canvas Name
        canvas_name_match = re.search(r"Canvas Name:\s*(.+)", response)
        if canvas_name_match:
            data["canvas_name"] = canvas_name_match.group(1).strip().strip("**")
        else:
            logger.warning("Canvas Name not found in the response.")

        # Extract Canvas Description
        canvas_description_match = re.search(r"Canvas Description:\s*(.+)", response)
        if canvas_description_match:
            data["canvas_description"] = canvas_description_match.group(1).strip().strip("**")
        else:
            logger.warning("Canvas Description not found in the response.")

        # Extract Hexagons Dynamically
        hexagon_sections = re.split(r"(?P<position>Top|Bottom) Hexagon \d+:", response)

        for i in range(1, len(hexagon_sections), 2):
            position = hexagon_sections[i].strip()
            content = hexagon_sections[i + 1].strip()

            # Extract Title, Description, and Key Elements
            title_match = re.search(r"Title:\s*(.+)", content)
            description_match = re.search(r"Description:\s*(.+)", content)
            key_elements_match = re.search(r"Key Elements:\s*(.+)", content)

            hexagon = {
                "title": title_match.group(1).strip() if title_match else None,
                "description": description_match.group(1).strip() if description_match else None,
                "key_elements": [
                    el.strip() for el in key_elements_match.group(1).split(",")
                ] if key_elements_match else [],
            }

            # Append to Top or Bottom Hexagons
            if position == "Top":
                data["top_hexagons"].append(hexagon)
            elif position == "Bottom":
                data["bottom_hexagons"].append(hexagon)

        # Log Missing Hexagons
        if not data["top_hexagons"]:
            logger.warning("'top_hexagons' is missing or empty.")
        if not data["bottom_hexagons"]:
            logger.warning("'bottom_hexagons' is missing or empty.")

        data = {key: clean_asterisks(value) if isinstance(value, str) else value for key, value in data.items()}
        for hex_list in ['top_hexagons', 'bottom_hexagons']:
            for hexagon in data[hex_list]:
                hexagon['title'] = clean_asterisks(hexagon['title'])
                hexagon['description'] = clean_asterisks(hexagon['description'])
                hexagon['key_elements'] = [clean_asterisks(k) for k in hexagon['key_elements']]

        return data
        
    except Exception as e:
        logger.error(f"Error parsing response: {str(e)}")
        raise ValueError(f"Parsing error: {str(e)}")
        
# Example functions to handle each template type
def handle_template_type_1(canvas_data):
    #presentation = Presentation("Progression Canvas.pptx")
    print(f"Handling template type 1 with data: {canvas_data}")
    presentation = Presentation("Hex Canvas Design (5).pptx")

    canvas_name = canvas_data.get("canvas_name", "")
    canvas_description = canvas_data.get("canvas_description", "")
    sections = canvas_data.get("sections", [])

    # Prepare replacement dictionary
    replacement_dict = {
        "cut1": canvas_name,
        "cut2": canvas_description,
    }

    # Dynamically map columns to placeholders (box1, box2, ..., box7)
    for section in sections:
        column = section.get("column", "")  # Extract column identifier
        if column.startswith("Column "):  # Ensure it matches the expected format
            column_number = column.split(" ")[1]  # Extract the column number
            placeholder = f"box{column_number}"  # Map to placeholders like box1, box2, etc.
            replacement_dict[placeholder] = {
                "title": section.get("title", "Default Title"),
                "description": section.get("description", "Default Description"),
                "key_elements": section.get("key_elements", [])[:4],  # Limit key elements to 4
            }

    # Debugging: Print the replacement dictionary
    print("==== TEMPLATE 1 REPLACEMENT DATA ====")
    for key, value in replacement_dict.items():
        if isinstance(value, dict):
            print(f"{key}: Title: {value['title']}, Description: {value['description']}, Key Elements: {', '.join(value['key_elements'])}")
        else:
            print(f"{key}: {value}")
    print("==== END OF REPLACEMENT DATA ====")

    # Function to apply replacements and formatting
    def apply_replacements(slide, replacement_dict):
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                for placeholder, data in replacement_dict.items():
                    if placeholder in shape.text:
                        # Replace placeholder with content
                        if isinstance(data, dict):
                            # Format for titles, descriptions, and key elements
                            formatted_text = (
                                f"{data['title']}\n\n{data['description']}\n- " + "\n- ".join(data["key_elements"])
                            )
                        else:
                            # Format for cut1 and cut2
                            formatted_text = data

                        shape.text = formatted_text  # Replace the text

                        # Apply formatting
                        if hasattr(shape, "text_frame") and shape.text_frame is not None:
                            shape.text_frame.word_wrap = True

                            for paragraph in shape.text_frame.paragraphs:
                                content = paragraph.text.strip()

                                # Format canvas name and description
                                if placeholder == "cut1":
                                    paragraph.alignment = PP_ALIGN.LEFT
                                    for run in paragraph.runs:
                                        run.font.bold = True
                                        run.font.size = Pt(20)
                                        run.font.color.rgb = RGBColor(0, 0, 0)  # Black
                                elif placeholder == "cut2":
                                    paragraph.alignment = PP_ALIGN.LEFT
                                    for run in paragraph.runs:
                                        run.font.size = Pt(16)
                                        run.font.color.rgb = RGBColor(0, 0, 0)  # Gray

                                # Format column titles
                                elif content == data.get("title", ""):
                                    paragraph.alignment = PP_ALIGN.CENTER
                                    for run in paragraph.runs:
                                        run.font.bold = True
                                        run.font.size = Pt(14)
                                        run.font.color.rgb = RGBColor(0, 0, 0)  # Black

                                # Format descriptions
                                elif content == data.get("description", ""):
                                    paragraph.alignment = PP_ALIGN.LEFT
                                    for run in paragraph.runs:
                                        run.font.size = Pt(12)
                                        run.font.color.rgb = RGBColor(0, 0, 0)  # Gray

                                # Format key elements
                                elif content.startswith("-"):
                                    paragraph.alignment = PP_ALIGN.LEFT
                                    paragraph.level = 1  # Indent for key elements
                                    for run in paragraph.runs:
                                        run.font.size = Pt(12)
                                        run.font.color.rgb = RGBColor(0, 0, 0)  # Gray

    # Apply replacements to all slides
    for slide in presentation.slides:
        apply_replacements(slide, replacement_dict)

    # Save the presentation and return as base64
    pptx_stream = BytesIO()
    presentation.save(pptx_stream)
    pptx_stream.seek(0)
    pptx_base64 = base64.b64encode(pptx_stream.read()).decode('utf-8')

    print("Template 1 processing complete.")
    return {"pptx_base64": pptx_base64}

def handle_template_type_2(canvas_data):
    print(f"Handling template type 2 with data: {canvas_data}")
    presentation = Presentation("Hex Canvas Design (7).pptx")  # Use the appropriate template for type 2

    canvas_name = canvas_data.get('canvas_name', '')
    canvas_description = canvas_data.get('canvas_description', '')
    sections = canvas_data.get('sections', [])
    replacement_dict = {}  # Dictionary to store all area data

    # Map canvas_name and description to cut1 and cut2
    replacement_dict['cut1'] = canvas_name
    replacement_dict['cut2'] = canvas_description

    # Map sections to their respective areas
    for section in sections:
        area = section.get('area', '')
        if area == "Top Right Area":
            replacement_dict["box1"] = {
                "title": section.get("title", "Default Title"),
                "description": section.get("description", "Default Description"),
                "key_elements": section.get("key_elements", [])[:3],  # Limit key elements to 3
            }
        elif area == "Top Left Area":
            replacement_dict["box2"] = {
                "title": section.get("title", "Default Title"),
                "description": section.get("description", "Default Description"),
                "key_elements": section.get("key_elements", [])[:3],
            }
        elif area == "Bottom Right Area":
            replacement_dict["box3"] = {
                "title": section.get("title", "Default Title"),
                "description": section.get("description", "Default Description"),
                "key_elements": section.get("key_elements", [])[:3],
            }
        elif area == "Bottom Left Area":
            replacement_dict["box4"] = {
                "title": section.get("title", "Default Title"),
                "description": section.get("description", "Default Description"),
                "key_elements": section.get("key_elements", [])[:3],
            }

    # Print the replacement_dict to verify its contents
    print("==== START OF BOX OUTPUT ====")
    for box, data in replacement_dict.items():
        print(f"{box} Data:")
        if isinstance(data, dict):
            print(f"Title: {data['title']}")
            if 'description' in data:
                print(f"Description: {data['description']}")
            if 'key_elements' in data:
                print(f"Key Elements: {', '.join(data['key_elements'])}")
        else:
            print(f"Content: {data}")
        print()
    print("==== END OF BOX OUTPUT ====")

    # Step 2: Apply replacements and handle formatting
    def apply_replacements(slide, replacement_dict, font_color_black):
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                for placeholder, data in replacement_dict.items():
                    if placeholder in shape.text:
                        # Replace placeholder with appropriate content
                        if isinstance(data, dict):
                            formatted_text = f"{data['title']}\n\n{data['description']}\n- " + "\n- ".join(data['key_elements'])
                        else:
                            formatted_text = data

                        shape.text = formatted_text  # Replace placeholder with text

                        # Apply formatting to the updated text
                        if hasattr(shape, "text_frame") and shape.text_frame is not None:
                            shape.text_frame.word_wrap = True

                            for paragraph in shape.text_frame.paragraphs:
                                # Extract the paragraph text
                                content = paragraph.text.strip()

                                # Apply formatting based on placeholders
                                if isinstance(data, dict):
                                    if content == data.get('title', ''):
                                        paragraph.alignment = PP_ALIGN.CENTER
                                        for run in paragraph.runs:
                                            run.font.bold = True
                                            run.font.size = Pt(14)  # Standard size for titles
                                            run.font.color.rgb = RGBColor(0, 0, 0) if font_color_black else RGBColor(255, 255, 255)
                                    elif content == data.get('description', ''):
                                        paragraph.alignment = PP_ALIGN.LEFT
                                        for run in paragraph.runs:
                                            run.font.size = Pt(12)
                                            run.font.color.rgb = RGBColor(0, 0, 0) if font_color_black else RGBColor(255, 255, 255)
                                    elif content.startswith("-"):
                                        paragraph.alignment = PP_ALIGN.LEFT
                                        paragraph.level = 1  # Indent for key elements
                                        for run in paragraph.runs:
                                            run.font.size = Pt(12)
                                            run.font.color.rgb = RGBColor(0, 0, 0) if font_color_black else RGBColor(255, 255, 255)

                                elif placeholder == "cut1":
                                    paragraph.alignment = PP_ALIGN.LEFT
                                    for run in paragraph.runs:
                                        run.font.bold = True
                                        run.font.size = Pt(20)
                                        run.font.color.rgb = RGBColor(0, 0, 0)  # Dark Blue
                                elif placeholder == "cut2":
                                    paragraph.alignment = PP_ALIGN.LEFT
                                    for run in paragraph.runs:
                                        run.font.size = Pt(16)
                                        run.font.color.rgb = RGBColor(0, 0, 0)  # Gray

    # Iterate through slides and apply replacements
    apply_replacements(presentation.slides[0], replacement_dict, font_color_black=True)  # Black font for Slide 1
    apply_replacements(presentation.slides[1], replacement_dict, font_color_black=False)  # White font for Slide 2

    # Save the presentation and return as base64
    pptx_stream = BytesIO()
    presentation.save(pptx_stream)
    pptx_stream.seek(0)
    pptx_base64 = base64.b64encode(pptx_stream.read()).decode('utf-8')

    print("Template 2 processing complete.")
    return {"pptx_base64": pptx_base64}
    
def handle_template_type_3(canvas_data):
    #presentation = Presentation("Circular Canvas.pptx")
    print(f"Handling template type 3 with data: {canvas_data}")
    presentation = Presentation("Hex Canvas Design (6).pptx")

    canvas_name = canvas_data.get("canvas_name", "")
    canvas_description = canvas_data.get("canvas_description", "")
    sections = canvas_data.get("sections", [])

    # Prepare the replacement dictionary
    replacement_dict_slide1 = {"cut1": canvas_name, "cut2": canvas_description}
    replacement_dict_slide2 = {"cut1": canvas_name, "cut2": canvas_description}

    box_counter = 1
    for section in sections:
        if section.get("circle") == "Central Circle":
            replacement_dict_slide1["center_circle"] = {
                "title": section.get("issue_goal", "Central Goal"),
            }
            replacement_dict_slide2["center_circle"] = replacement_dict_slide1["center_circle"]

        elif section.get("circle", "").startswith("Supporting Circle"):
            raw_key_elements = section.get("key_elements", [])
            cleaned_key_elements = []

            # Clean up key elements and stop parsing if new section starts
            for element in raw_key_elements:
                if "Supporting Circle" in element:
                    break
                cleaned_key_elements.append(element.strip())

            # Add to replacement dictionaries
            replacement_dict_slide1[f"box{box_counter}"] = {"title": section.get("title", "Default Title")}
            replacement_dict_slide2[f"box{box_counter}"] = {
                "title": section.get("title", "Default Title"),
                "description": section.get("description", "Default Description"),
                "key_elements": cleaned_key_elements[:3],
            }
            box_counter += 1

    # Print the replacement dictionaries for debugging
    print("==== SLIDE 1 REPLACEMENT DATA ====")
    print(replacement_dict_slide1)
    print("==== SLIDE 2 REPLACEMENT DATA ====")
    print(replacement_dict_slide2)

    # Function to apply replacements and formatting
    def apply_replacements(slide, replacement_dict, font_color_black=True):
        """
        Applies replacements for placeholders in the provided slide.
        Handles formatting dynamically based on the placeholder type.
        """
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                for placeholder, data in replacement_dict.items():
                    if placeholder in shape.text:
                        # Replace placeholder with content
                        if isinstance(data, dict):
                            if "description" in data and "key_elements" in data:
                                formatted_text = (
                                    f"{data['title']}\n\n{data['description']}\n- " + "\n- ".join(data["key_elements"])
                                )
                            else:
                                formatted_text = data["title"]
                        else:
                            formatted_text = data

                        # Replace the placeholder text
                        shape.text = formatted_text

                        # Apply formatting to the updated text
                        if hasattr(shape, "text_frame") and shape.text_frame is not None:
                            for paragraph in shape.text_frame.paragraphs:
                                content = paragraph.text.strip()

                                # Formatting for canvas name and description
                                if placeholder == "cut1":
                                    paragraph.alignment = PP_ALIGN.LEFT
                                    for run in paragraph.runs:
                                        run.font.bold = True
                                        run.font.size = Pt(20)
                                        run.font.color.rgb = RGBColor(0, 0, 0)  # Black
                                elif placeholder == "cut2":
                                    paragraph.alignment = PP_ALIGN.LEFT
                                    for run in paragraph.runs:
                                        run.font.size = Pt(16)
                                        run.font.color.rgb = RGBColor(0, 0, 0)  # Black

                                # Formatting for titles
                                elif content == data.get("title", ""):
                                    if placeholder == "center_circle":
                                        paragraph.alignment = PP_ALIGN.CENTER
                                        for run in paragraph.runs:
                                            #run.font.bold = True
                                            run.font.size = Pt(14)  # Larger font for center circle
                                            run.font.color.rgb = RGBColor(255, 255, 255)  # Red
                                    else:
                                        paragraph.alignment = PP_ALIGN.CENTER
                                        for run in paragraph.runs:
                                            run.font.bold = True
                                            run.font.size = Pt(14)
                                            run.font.color.rgb = (
                                                RGBColor(0, 0, 0) if font_color_black else RGBColor(255, 255, 255)
                                            )

                                # Formatting for descriptions
                                elif content == data.get("description", ""):
                                    paragraph.alignment = PP_ALIGN.LEFT
                                    for run in paragraph.runs:
                                        run.font.size = Pt(12)
                                        run.font.color.rgb = (
                                            RGBColor(0, 0, 0) if font_color_black else RGBColor(255, 255, 255)
                                        )

                                # Formatting for key elements
                                elif content.startswith("-"):
                                    paragraph.alignment = PP_ALIGN.LEFT
                                    paragraph.level = 1  # Indentation for key elements
                                    for run in paragraph.runs:
                                        run.font.size = Pt(12)
                                        run.font.color.rgb = (
                                            RGBColor(0, 0, 0) if font_color_black else RGBColor(255, 255, 255)
                                        )

    # Apply replacements for Slide 1 (titles only, black font)
    apply_replacements(presentation.slides[0], replacement_dict_slide2, font_color_black=True)

    # Apply replacements for Slide 2 (titles, descriptions, key elements, white font)
    apply_replacements(presentation.slides[1], replacement_dict_slide1, font_color_black=True)

    # Save the presentation and return as base64
    pptx_stream = BytesIO()
    presentation.save(pptx_stream)
    pptx_stream.seek(0)
    pptx_base64 = base64.b64encode(pptx_stream.read()).decode("utf-8")

    print("Template 3 processing complete.")
    return {"pptx_base64": pptx_base64}
    
def handle_template_type_4(canvas_data):
    presentation = Presentation("Hex Canvas Design (3).pptx")
    print(f"Handling template type 4 with data: {canvas_data}")
    if "top_hexagons" not in canvas_data or not canvas_data["top_hexagons"]:
        raise ValueError("'top_hexagons' is missing or empty.")
    if "bottom_hexagons" not in canvas_data or not canvas_data["bottom_hexagons"]:
        raise ValueError("'bottom_hexagons' is missing or empty.")

    print("Top hexagons:", canvas_data["top_hexagons"])
    print("Bottom hexagons:", canvas_data["bottom_hexagons"])
    # Adjust the replacement dictionary, including 'cut1' and 'cut2' with different font sizes and title color
    # Build the complete replacement dictionary to handle titles, descriptions, and key elements
    replacement_dict_slide1 = {
        "box1": canvas_data['top_hexagons'][0]['title'],
        "top_hex2": canvas_data['top_hexagons'][1]['title'],
        "top_hex3": canvas_data['top_hexagons'][2]['title'],
        "top_hex4": canvas_data['top_hexagons'][3]['title'],
        "box2": canvas_data['bottom_hexagons'][0]['title'],
        "bottom_hex2": canvas_data['bottom_hexagons'][1]['title'],
        "bottom_hex3": canvas_data['bottom_hexagons'][2]['title'],
        "bottom_hex4": canvas_data['bottom_hexagons'][3]['title'],
        "cut1": canvas_data["canvas_name"],
        "cut2": canvas_data["canvas_description"],
    }
    
    replacement_dict_slide2 = {
        "box1": (
            f"{canvas_data['top_hexagons'][0]['title']}\n\n"
            f"{canvas_data['top_hexagons'][0]['description']}\n- "
            + "\n- ".join(canvas_data['top_hexagons'][0]['key_elements'][:4])
        ),
        "top_hex2": (
            f"{canvas_data['top_hexagons'][1]['title']}\n\n"
            f"{canvas_data['top_hexagons'][1]['description']}\n- "
            + "\n- ".join(canvas_data['top_hexagons'][1]['key_elements'][:4])
        ),
        "top_hex3": (
            f"{canvas_data['top_hexagons'][2]['title']}\n\n"
            f"{canvas_data['top_hexagons'][2]['description']}\n- "
            + "\n- ".join(canvas_data['top_hexagons'][2]['key_elements'][:4])
        ),
        "top_hex4": (
            f"{canvas_data['top_hexagons'][3]['title']}\n\n"
            f"{canvas_data['top_hexagons'][3]['description']}\n- "
            + "\n- ".join(canvas_data['top_hexagons'][3]['key_elements'][:4])
        ),
        "box2": (
            f"{canvas_data['bottom_hexagons'][0]['title']}\n\n"
            f"{canvas_data['bottom_hexagons'][0]['description']}\n- "
            + "\n- ".join(canvas_data['bottom_hexagons'][0]['key_elements'][:4])
        ),
        "bottom_hex2": (
            f"{canvas_data['bottom_hexagons'][1]['title']}\n\n"
            f"{canvas_data['bottom_hexagons'][1]['description']}\n- "
            + "\n- ".join(canvas_data['bottom_hexagons'][1]['key_elements'][:4])
        ),
        "bottom_hex3": (
            f"{canvas_data['bottom_hexagons'][2]['title']}\n\n"
            f"{canvas_data['bottom_hexagons'][2]['description']}\n- "
            + "\n- ".join(canvas_data['bottom_hexagons'][2]['key_elements'][:4])
        ),
        "bottom_hex4": (
            f"{canvas_data['bottom_hexagons'][3]['title']}\n\n"
            f"{canvas_data['bottom_hexagons'][3]['description']}\n- "
            + "\n- ".join(canvas_data['bottom_hexagons'][3]['key_elements'][:4])
        ),
        "cut1": canvas_data["canvas_name"],
        "cut2": canvas_data["canvas_description"],
    }
    
    
    # Iterate through slides and apply formatting for 'cut1', 'cut2', hexagon titles, descriptions, and key elements
    def apply_replacements(slide, replacement_dict, font_color_black=False):
        """
        Applies replacements for placeholders in the provided slide.
        Adjusts formatting based on the placeholder type.
        """
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                print(f"Checking shape text: {shape.text}")
                for placeholder, replacement in replacement_dict.items():
                    if placeholder in shape.text:
                        print(f"Replacing placeholder '{placeholder}' with '{replacement}'")
                        shape.text = shape.text.replace(placeholder, replacement)
    
                        # Apply text formatting
                        if hasattr(shape, "text_frame") and shape.text_frame is not None:
                            shape.text_frame.word_wrap = True
                            shape.text_frame.auto_size = True
    
                            for paragraph in shape.text_frame.paragraphs:
                                text_content = paragraph.text.strip()
    
                                # Apply specific styles for cut1 and cut2
                                if placeholder == "cut1":
                                    for run in paragraph.runs:
                                        run.font.size = Pt(20)
                                        run.font.bold = True
                                        run.font.name = "Arial"
                                        run.font.color.rgb = RGBColor(0, 0, 0)  # Black
                                    continue
    
                                elif placeholder == "cut2":
                                    for run in paragraph.runs:
                                        run.font.size = Pt(14)
                                        run.font.bold = True
                                        run.font.name = "Arial"
                                        run.font.color.rgb = RGBColor(0, 0, 0)  # Black
                                    continue
    
                                # Description formatting
                                if replacement in text_content:
                                    paragraph.level = 0
                                    paragraph.alignment = PP_ALIGN.LEFT
                                    for run in paragraph.runs:
                                        run.font.size = Pt(12)
                                        run.font.color.rgb = RGBColor(255, 255, 255)
                                    continue
    
                                # Title formatting: Center-align and make bold
                                if replacement in replacement_dict.values():
                                    paragraph.alignment = PP_ALIGN.CENTER
                                    for run in paragraph.runs:
                                        run.font.bold = True
                                        run.font.size = Pt(14)
                                        run.font.color.rgb = RGBColor(0, 0, 0) if font_color_black else RGBColor(255, 255, 255)
                                    continue
    
                                # Default formatting for unhandled cases
                                for run in paragraph.runs:
                                    run.font.size = Pt(11)
                                    run.font.name = "Arial"
                                    run.font.color.rgb = RGBColor(255, 255, 255)  # White

    apply_replacements(presentation.slides[0], replacement_dict_slide1, font_color_black=True)  # Slide 1: Titles only, black font
    apply_replacements(presentation.slides[1], replacement_dict_slide2)  # Slide 2: Titles, descriptions, key elements

    pptx_stream = BytesIO()
    presentation.save(pptx_stream)
    pptx_stream.seek(0)  # Move the stream position to the start
    pptx_base64 = base64.b64encode(pptx_stream.read()).decode('utf-8')
    
    log_memory_usage("Save Presentation")
    print("done")
    return {
        'pptx_base64': pptx_base64,
    }
    
    
    



    
